from __future__ import annotations

import csv
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
CSV_PATH = ROOT / "data" / "brand_json_summaries" / "ttamtti2_local_review_analysis.csv"
OUT_DIR = ROOT / "outputs" / "cosmetic_example"
OUT_PPTX = OUT_DIR / "three_brand_review_summary_sample.pptx"

BRANDS = ["1.618", "HYCL", "헤이네이처"]

BLACK = RGBColor(16, 16, 16)
MUTED = RGBColor(88, 88, 88)
LIGHT = RGBColor(242, 242, 242)
MID = RGBColor(214, 218, 224)
ORANGE = RGBColor(255, 107, 53)
TEAL = RGBColor(20, 126, 116)
BLUE = RGBColor(42, 91, 215)
WHITE = RGBColor(255, 255, 255)
PALE_ORANGE = RGBColor(255, 244, 239)
PALE_TEAL = RGBColor(235, 247, 245)


def read_rows() -> dict[str, dict[str, str]]:
    with CSV_PATH.open("r", encoding="utf-8-sig", newline="") as f:
        rows = list(csv.DictReader(f))
    return {row["브랜드명"]: row for row in rows}


def add_text(slide, text, x, y, w, h, size=18, color=BLACK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def rect(slide, x, y, w, h, fill=LIGHT, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.58, 0.42, 11.6, 0.62, 28, BLACK, True)
    rect(slide, 0.61, 1.23, 11.75, 0.015, MID)
    if subtitle:
        add_text(slide, subtitle, 0.61, 1.35, 11.1, 0.32, 12, MUTED)


def add_footer(slide, page):
    add_text(slide, "브랜드별 리뷰 정리 샘플", 0.61, 7.08, 3.0, 0.24, 8, MUTED)
    add_text(slide, f"{page:02d}", 12.0, 7.08, 0.5, 0.24, 8, MUTED, align=PP_ALIGN.RIGHT)


def split_items(value: str, limit: int = 3) -> list[str]:
    parts = [p.strip() for p in (value or "").split(" / ") if p.strip()]
    return parts[:limit]


def clean_review(text: str, limit: int = 84) -> str:
    text = re.sub(r"\([^)]*등록된 네이버 페이 구매평\)", "", text or "")
    text = re.sub(r"\d{4}[-.]\d{1,2}[-.]\d{1,2}", "", text)
    text = re.sub(r"\b\d{1,2}:\d{2}:\d{2}\b", "", text)
    text = re.sub(r"[A-Za-z]{2,}\*{2,}", "", text)
    text = re.sub(r"네\*{2,}|[가-힣A-Za-z]\*{2,}", "", text)
    text = " ".join(text.replace("더보기", "").split())
    if len(text) > limit:
        text = text[: limit - 1].rstrip() + "..."
    return text


def keywords(row: dict[str, str], limit: int = 6) -> list[str]:
    bad = {"TO", "SHIPPING", "ABOUT", "REVIEW", "BEST", "EVENT", "Login", "Cart"}
    out = []
    for item in (row.get("상위키워드") or "").split(","):
        word = item.strip()
        if not word or word in bad:
            continue
        if re.search(r"[가-힣]", word) or len(word) <= 12:
            out.append(word)
        if len(out) >= limit:
            break
    return out


def bullet_list(slide, items, x, y, w, h, size=12, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return box


def tag(slide, text, x, y, w, color):
    rect(slide, x, y, w, 0.34, color)
    add_text(slide, text, x + 0.12, y + 0.07, w - 0.24, 0.16, 9, WHITE, True, PP_ALIGN.CENTER)


def metric(slide, label, value, x, y, w, color=BLACK):
    rect(slide, x, y, w, 0.84, LIGHT)
    add_text(slide, value, x + 0.14, y + 0.13, w - 0.28, 0.28, 18, color, True)
    add_text(slide, label, x + 0.14, y + 0.51, w - 0.28, 0.18, 8.5, MUTED)


def add_brand_slide(prs: Presentation, row: dict[str, str], page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    brand = row["브랜드명"]
    add_title(slide, f"{brand} 리뷰 정리", "리뷰에서 반복되는 고객 언어를 영업 제안 포인트로 변환했습니다.")

    tag(slide, row["대표그룹"], 0.7, 1.85, 2.3, TEAL)
    tag(slide, row["세부유형"], 3.2, 1.85, 2.2, ORANGE)
    tag(slide, f"근거 {row['근거등급']}등급", 5.6, 1.85, 1.4, BLUE)

    metric(slide, "정제 리뷰 문장", row["정제문장수"], 0.72, 2.48, 1.85, BLACK)
    metric(slide, "긍정 문장", row["긍정문장수"], 2.78, 2.48, 1.85, TEAL)
    metric(slide, "부정/주의 문장", row["부정문장수"], 4.84, 2.48, 1.85, ORANGE)

    add_text(slide, "핵심 키워드", 7.25, 2.45, 2.0, 0.26, 15, BLACK, True)
    for i, kw in enumerate(keywords(row, 6)):
        x = 7.25 + (i % 3) * 1.38
        y = 2.88 + (i // 3) * 0.43
        rect(slide, x, y, 1.18, 0.29, PALE_TEAL, TEAL)
        add_text(slide, kw, x + 0.06, y + 0.06, 1.06, 0.13, 8.2, TEAL, True, PP_ALIGN.CENTER)

    add_text(slide, "리뷰 근거", 0.72, 3.78, 2.0, 0.28, 16, BLACK, True)
    review_items = [clean_review(s, 68) for s in split_items(row["대표문장"], 3)]
    bullet_list(slide, review_items, 0.78, 4.22, 5.55, 1.45, 11.5)

    add_text(slide, "영업 제안", 6.95, 3.78, 2.0, 0.28, 16, BLACK, True)
    recs = split_items(row["추천영업상품"], 4)
    bullet_list(slide, recs, 7.02, 4.22, 4.8, 1.1, 12.5)

    rect(slide, 6.98, 5.65, 5.1, 0.72, PALE_ORANGE)
    add_text(slide, "제안 메시지", 7.18, 5.82, 1.35, 0.2, 10, ORANGE, True)
    add_text(slide, clean_review(row["제안브리프"], 58), 8.42, 5.76, 3.35, 0.32, 11.5, BLACK, True)

    add_footer(slide, page)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    rows_by_brand = read_rows()
    selected = [rows_by_brand[name] for name in BRANDS if name in rows_by_brand]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(cover, "브랜드별 리뷰 정리\n3개 샘플", 0.66, 0.9, 7.2, 1.4, 36, BLACK, True)
    add_text(cover, "로컬 JSON 분석 결과를 바탕으로 브랜드별 리뷰 핵심, 고객 신호, 추천 영업상품을 한 장씩 정리했습니다.", 0.7, 2.65, 7.6, 0.48, 15, MUTED)
    rect(cover, 8.35, 0.86, 4.1, 5.55, LIGHT)
    add_text(cover, "포함 브랜드", 8.7, 1.18, 2.2, 0.28, 17, BLACK, True)
    for idx, row in enumerate(selected, start=1):
        y = 1.78 + (idx - 1) * 1.15
        add_text(cover, f"0{idx}", 8.72, y, 0.55, 0.34, 18, ORANGE, True)
        add_text(cover, row["브랜드명"], 9.35, y, 2.5, 0.34, 18, BLACK, True)
        add_text(cover, row["대표그룹"], 9.35, y + 0.38, 2.7, 0.22, 10, MUTED)
    add_footer(cover, 1)

    for page, row in enumerate(selected, start=2):
        add_brand_slide(prs, row, page)

    close = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(close, "3개 브랜드 모두 제안상품을 다르게 가져갈 수 있다", "리뷰 정리는 단순 요약이 아니라, 브랜드마다 다른 영업 제안의 출발점입니다.")
    summary = [
        ("1.618", "촉촉함·재구매 신호", "후기형 상세페이지와 혜택 CRM"),
        ("HYCL", "순한 사용감·클렌징/세럼 리뷰", "사용감 숏폼과 검색광고 키워드"),
        ("헤이네이처", "어성초·진정·수분장벽", "민감 피부 근거 강화 상세페이지"),
    ]
    for i, (brand, signal, action) in enumerate(summary):
        y = 2.15 + i * 1.1
        rect(close, 0.82, y, 11.4, 0.78, LIGHT)
        add_text(close, brand, 1.08, y + 0.19, 1.65, 0.24, 16, BLACK, True)
        add_text(close, signal, 3.1, y + 0.2, 3.05, 0.22, 13, TEAL, True)
        add_text(close, action, 6.55, y + 0.2, 4.6, 0.22, 13, BLACK)
    rect(close, 0.82, 5.95, 11.4, 0.52, PALE_ORANGE)
    add_text(close, "다음에는 이 구조를 300개 브랜드에 반복 적용해, 브랜드별 PPT와 제안메일 본문을 함께 뽑으면 됩니다.", 1.08, 6.09, 10.6, 0.22, 13, BLACK, True)
    add_footer(close, len(selected) + 2)

    prs.save(OUT_PPTX)
    print(OUT_PPTX)


if __name__ == "__main__":
    main()
