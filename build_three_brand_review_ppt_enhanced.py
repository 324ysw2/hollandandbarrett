from __future__ import annotations

import csv
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DATA_DIR = ROOT / "data" / "brand_json_summaries"
LOCAL_CSV = DATA_DIR / "ttamtti2_local_review_analysis.csv"
BRIEF_CSV = DATA_DIR / "ttamtti2_sales_automation_briefs.csv"
OUT_DIR = ROOT / "outputs" / "cosmetic_example"
OUT_PPTX = OUT_DIR / "three_brand_review_summary_enhanced.pptx"

BRANDS = ["1.618", "HYCL", "헤이네이처"]

BLACK = RGBColor(16, 16, 16)
MUTED = RGBColor(86, 86, 86)
LIGHT = RGBColor(244, 244, 244)
MID = RGBColor(214, 218, 224)
ORANGE = RGBColor(255, 107, 53)
TEAL = RGBColor(20, 126, 116)
BLUE = RGBColor(42, 91, 215)
WHITE = RGBColor(255, 255, 255)
PALE_ORANGE = RGBColor(255, 244, 239)
PALE_TEAL = RGBColor(235, 247, 245)


BAD_REVIEW_TERMS = {
    "ABOUT",
    "EVENT",
    "BEST",
    "SHOP",
    "COMMUNITY",
    "Login",
    "Join",
    "Cart",
    "자세히보기",
    "검색 가입",
    "브랜드 스토리",
    "Social Contribution",
    "Clean Beauty",
    "글로우픽 앱에서 보기",
}


def read_csv_map(path: Path) -> dict[str, dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return {row["브랜드명"]: row for row in csv.DictReader(f)}


def add_text(slide, text, x, y, w, h, size=18, color=BLACK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def rect(slide, x, y, w, h, fill=LIGHT, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.58, 0.42, 11.6, 0.62, 28, BLACK, True)
    rect(slide, 0.61, 1.23, 11.75, 0.015, MID)
    if subtitle:
        add_text(slide, subtitle, 0.61, 1.35, 11.2, 0.32, 12, MUTED)


def add_footer(slide, page):
    add_text(slide, "브랜드별 리뷰 정리", 0.61, 7.08, 3.0, 0.24, 8, MUTED)
    add_text(slide, f"{page:02d}", 12.0, 7.08, 0.5, 0.24, 8, MUTED, align=PP_ALIGN.RIGHT)


def split_slash(value: str, limit: int | None = None) -> list[str]:
    items = [" ".join(v.strip().split()) for v in (value or "").split(" / ") if v.strip()]
    return items[:limit] if limit else items


def split_comma(value: str, limit: int | None = None) -> list[str]:
    items = [" ".join(v.strip().split()) for v in (value or "").split(",") if v.strip()]
    return items[:limit] if limit else items


def clean_text(text: str) -> str:
    text = re.sub(r"\([^)]*등록된 네이버 페이 구매평\)", " ", text or "")
    text = re.sub(r"\d{4}[-.]\d{1,2}[-.]\d{1,2}(?:\s+\d{1,2}:\d{2}:\d{2})?", " ", text)
    text = re.sub(r"\b\d{1,2}:\d{2}:\d{2}\b", " ", text)
    text = re.sub(r"[A-Za-z가-힣]\*{2,}", " ", text)
    text = re.sub(r"\b\d{1,3}(?:,\d{3})*원\b", " ", text)
    text = text.replace("더보기", " ").replace("[수정됨]", " ")
    return " ".join(text.split())


def sentence_parts(value: str) -> list[str]:
    parts = []
    for chunk in split_slash(value):
        chunk = clean_text(chunk)
        parts.extend(re.split(r"(?<=[.!?。！？])\s+|(?<=요)\s+|(?<=다)\s+", chunk))
    return [" ".join(part.split()) for part in parts if part.strip()]


def review_score(text: str) -> int:
    score = 0
    good_terms = [
        "좋",
        "촉촉",
        "진정",
        "순하",
        "자극",
        "부드럽",
        "재구매",
        "만족",
        "세안",
        "발림",
        "피부",
        "여드",
        "수분",
        "장벽",
        "흡수",
    ]
    score += sum(2 for term in good_terms if term in text)
    if 18 <= len(text) <= 86:
        score += 8
    if len(text) > 110:
        score -= 8
    if any(term in text for term in BAD_REVIEW_TERMS):
        score -= 15
    if text.count(" ") < 2:
        score -= 4
    return score


def best_reviews(local: dict[str, str], limit: int = 3) -> list[str]:
    candidates = []
    for field in ["긍정예시", "대표문장", "부정예시"]:
        candidates.extend(sentence_parts(local.get(field, "")))
    unique = []
    seen = set()
    for item in candidates:
        if len(item) < 12 or item in seen:
            continue
        seen.add(item)
        unique.append(item)
    unique.sort(key=review_score, reverse=True)
    picked = []
    for item in unique:
        if review_score(item) < 2:
            continue
        if len(item) > 92:
            continue
        picked.append(item)
        if len(picked) >= limit:
            break
    return picked or [clean_text(p)[:86] for p in sentence_parts(local.get("대표문장", ""))[:limit]]


def tag(slide, text, x, y, w, color):
    rect(slide, x, y, w, 0.34, color)
    add_text(slide, text, x + 0.08, y + 0.07, w - 0.16, 0.15, 8.6, WHITE, True, PP_ALIGN.CENTER)


def metric(slide, label, value, x, y, w, color=BLACK):
    rect(slide, x, y, w, 0.72, LIGHT)
    add_text(slide, value, x + 0.12, y + 0.11, w - 0.24, 0.24, 16, color, True)
    add_text(slide, label, x + 0.12, y + 0.43, w - 0.24, 0.16, 8.2, MUTED)


def bullet_list(slide, items, x, y, w, h, size=11.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(6)
    return box


def add_brand_slide(prs: Presentation, local: dict[str, str], brief: dict[str, str], page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    brand = local["브랜드명"]
    add_title(slide, f"{brand} 리뷰 정리", "리뷰 근거와 추천상품을 한 장에서 바로 확인합니다.")

    tag(slide, local["대표그룹"], 0.7, 1.82, 2.35, TEAL)
    tag(slide, local["세부유형"], 3.22, 1.82, 1.95, ORANGE)
    tag(slide, f"근거 {local['근거등급']}등급", 5.35, 1.82, 1.22, BLUE)

    metric(slide, "리뷰 문장", local["정제문장수"], 0.72, 2.35, 1.55, BLACK)
    metric(slide, "긍정", local["긍정문장수"], 2.46, 2.35, 1.35, TEAL)
    metric(slide, "주의", local["부정문장수"], 3.98, 2.35, 1.35, ORANGE)
    metric(slide, "우선순위", brief["우선순위점수"], 5.5, 2.35, 1.45, BLUE)

    add_text(slide, "리뷰 근거", 0.72, 3.42, 1.55, 0.28, 16, BLACK, True)
    bullet_list(slide, best_reviews(local, 3), 0.78, 3.82, 5.55, 1.32, 11.4)

    add_text(slide, "구매장벽", 0.72, 5.52, 1.55, 0.26, 15, BLACK, True)
    bullet_list(slide, split_slash(brief["구매장벽"], 3), 0.78, 5.88, 5.55, 0.65, 11)

    add_text(slide, "추천상품", 6.78, 3.42, 1.55, 0.28, 16, BLACK, True)
    bullet_list(slide, split_slash(brief["추천영업상품"], 4), 6.85, 3.82, 4.95, 1.16, 11.8)

    add_text(slide, "유사 브랜드", 6.78, 5.28, 1.55, 0.26, 15, BLACK, True)
    similar = split_slash(brief["유사브랜드"], 4)
    add_text(slide, " / ".join(similar), 6.85, 5.64, 4.95, 0.26, 11.2, TEAL, True)

    rect(slide, 6.82, 6.05, 4.95, 0.52, PALE_ORANGE)
    add_text(slide, "제안 문장", 7.02, 6.18, 1.0, 0.16, 9.4, ORANGE, True)
    one_line = f"{brand}는 {', '.join(split_comma(brief['공유키워드'], 3))} 신호가 있어 {split_slash(brief['추천영업상품'], 1)[0]} 제안이 적합합니다."
    add_text(slide, one_line, 8.0, 6.13, 3.45, 0.22, 10.4, BLACK, True)

    add_footer(slide, page)


def add_cover(prs: Presentation, locals_by_brand: dict[str, dict[str, str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, "브랜드별 리뷰 정리\n고도화 샘플", 0.66, 0.88, 7.4, 1.4, 36, BLACK, True)
    add_text(slide, "리뷰 근거, 구매장벽, 유사 브랜드, 추천상품만 남긴 제안서용 구성입니다.", 0.7, 2.62, 7.2, 0.46, 15, MUTED)
    rect(slide, 8.35, 0.86, 4.1, 5.55, LIGHT)
    add_text(slide, "포함 브랜드", 8.7, 1.18, 2.2, 0.28, 17, BLACK, True)
    for idx, brand in enumerate(BRANDS, start=1):
        row = locals_by_brand[brand]
        y = 1.78 + (idx - 1) * 1.15
        add_text(slide, f"0{idx}", 8.72, y, 0.55, 0.34, 18, ORANGE, True)
        add_text(slide, row["브랜드명"], 9.35, y, 2.5, 0.34, 18, BLACK, True)
        add_text(slide, row["대표그룹"], 9.35, y + 0.38, 2.7, 0.22, 10, MUTED)
    add_footer(slide, 1)


def add_close(prs: Presentation, briefs_by_brand: dict[str, dict[str, str]], page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "리뷰가 말하는 문제에 맞춰 제안상품을 다르게 잡습니다", "같은 화장품 브랜드라도 반복되는 리뷰 신호가 다르면 영업 제안도 달라집니다.")
    for idx, brand in enumerate(BRANDS):
        brief = briefs_by_brand[brand]
        y = 2.0 + idx * 1.16
        rect(slide, 0.82, y, 11.4, 0.84, LIGHT)
        add_text(slide, brand, 1.08, y + 0.2, 1.65, 0.24, 16, BLACK, True)
        add_text(slide, split_slash(brief["구매장벽"], 1)[0], 3.05, y + 0.2, 3.2, 0.22, 12.2, TEAL, True)
        add_text(slide, split_slash(brief["추천영업상품"], 1)[0], 6.58, y + 0.2, 4.6, 0.22, 12.2, BLACK)
    rect(slide, 0.82, 5.78, 11.4, 0.62, PALE_ORANGE)
    add_text(slide, "다음 작업", 1.08, 5.98, 1.1, 0.18, 10.5, ORANGE, True)
    add_text(slide, "이 구조를 브랜드별 자동 PPT, 이미지 카드, 제안메일 본문에 그대로 반복 적용합니다.", 2.18, 5.9, 8.7, 0.28, 14, BLACK, True)
    add_footer(slide, page)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    locals_by_brand = read_csv_map(LOCAL_CSV)
    briefs_by_brand = read_csv_map(BRIEF_CSV)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs, locals_by_brand)
    for page, brand in enumerate(BRANDS, start=2):
        add_brand_slide(prs, locals_by_brand[brand], briefs_by_brand[brand], page)
    add_close(prs, briefs_by_brand, len(BRANDS) + 2)

    prs.save(OUT_PPTX)
    print(OUT_PPTX)


if __name__ == "__main__":
    main()
