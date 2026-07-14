from __future__ import annotations

import csv
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
CSV = ROOT / "data" / "brand_json_summaries" / "ttamtti2_cosmetic_sales_grouping.csv"
OUT_DIR = ROOT / "outputs" / "cosmetic_example"
IMAGE_OUT = OUT_DIR / "goongbe_email_example.png"
PPT_OUT = OUT_DIR / "goongbe_sales_proposal_example.pptx"

WIDE_W, WIDE_H = 13.333, 7.5


def font_path(bold: bool = False) -> str:
    candidates = [
        r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf",
        r"C:\Windows\Fonts\NotoSansKR-Bold.ttf" if bold else r"C:\Windows\Fonts\NotoSansKR-Regular.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return path
    return ""


def pil_font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    path = font_path(bold)
    if path:
        return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def wrap_kr(text: str, width: int) -> str:
    return "\n".join(textwrap.wrap(str(text), width=width, break_long_words=False, replace_whitespace=False))


def read_brand(brand_name: str = "궁중비책") -> dict[str, str]:
    with CSV.open("r", encoding="utf-8-sig", newline="") as handle:
        for row in csv.DictReader(handle):
            if row.get("브랜드명") == brand_name:
                return row
    raise RuntimeError(f"Brand not found: {brand_name}")


def rounded(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], radius: int, fill, outline=None, width: int = 1) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def make_email_image(row: dict[str, str]) -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (1200, 675), "#F6F3EE")
    draw = ImageDraw.Draw(img)

    green = "#1F6F5B"
    deep = "#18201E"
    sage = "#CFE2D8"
    mint = "#E7F2EC"
    cream = "#FFF8ED"
    coral = "#E47C64"
    ink = "#202624"
    muted = "#68736F"

    draw.rectangle((0, 0, 1200, 675), fill="#F6F3EE")
    draw.ellipse((760, -160, 1320, 440), fill="#DCEBE3")
    draw.ellipse((830, 150, 1235, 760), fill="#F1DDCB")
    draw.ellipse((700, 350, 955, 610), fill="#BFD8CD")

    # Product-inspired soft bottle and cream jar silhouettes.
    rounded(draw, (805, 110, 950, 500), 44, fill="#FFFFFF", outline="#D6DFDA", width=3)
    rounded(draw, (835, 72, 920, 132), 18, fill="#DDE8E1", outline="#B9CCC3", width=2)
    rounded(draw, (985, 305, 1130, 515), 38, fill="#FFFFFF", outline="#D6DFDA", width=3)
    rounded(draw, (968, 270, 1147, 325), 24, fill="#E5EEE9", outline="#C1D2CA", width=2)
    draw.text((841, 240), "GOONGBE", font=pil_font(25, True), fill=green)
    draw.text((1015, 394), "Barrier\nCream", font=pil_font(25, True), fill=green, spacing=4)

    rounded(draw, (62, 58, 244, 98), 20, fill=green)
    draw.text((88, 68), "제안메일 예시", font=pil_font(22, True), fill="#FFFFFF")

    draw.text((66, 132), "궁중비책", font=pil_font(70, True), fill=deep)
    draw.text((70, 218), "진정과 장벽 케어를\n후기형 콘텐츠로 설득합니다", font=pil_font(43, True), fill=ink, spacing=8)

    body = "리뷰/FAQ 분석상 진정, 장벽, 시카, 병풀 신호가 반복됩니다. 상세페이지와 메일 이미지는 효능 단정보다 사용감과 구매 전 확인 포인트를 중심으로 정리하는 편이 안전합니다."
    draw.text((72, 335), wrap_kr(body, 34), font=pil_font(24), fill=muted, spacing=8)

    badges = [
        ("진정/장벽", green),
        ("후기 기반 만족", coral),
        ("구매 혜택 강조", "#C59043"),
    ]
    x = 72
    for label, color in badges:
        tw = int(draw.textlength(label, font=pil_font(22, True)))
        rounded(draw, (x, 520, x + tw + 42, 568), 24, fill=color)
        draw.text((x + 21, 532), label, font=pil_font(22, True), fill="#FFFFFF")
        x += tw + 58

    rounded(draw, (72, 590, 655, 635), 12, fill=cream, outline="#E1D7C9", width=2)
    draw.text((96, 602), "추천 영업상품: 체험단/숏폼 리뷰 콘텐츠 + 상세페이지 개선", font=pil_font(20, True), fill=green)
    draw.text((825, 540), "PPT/이미지\n제안용 샘플", font=pil_font(31, True), fill=deep, spacing=8)
    img.save(IMAGE_OUT, quality=95)


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_text(slide, text, x, y, w, h, size=18, color="#202624", bold=False, align="left", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    return box


def add_rect(slide, x, y, w, h, fill="#FFFFFF", line="#D8DED9", radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    return shape


def add_badge(slide, text, x, y, color="#1F6F5B", w=None):
    width = w if w is not None else max(1.25, min(2.4, 0.36 + len(text) * 0.13))
    add_rect(slide, x, y, width, 0.36, fill=color, line=color, radius=True)
    add_text(slide, text, x + 0.07, y + 0.08, width - 0.14, 0.14, 10, "#FFFFFF", True, "center")


def add_footer(slide, page: int):
    add_text(slide, "Cosmetic sales proposal sample", 0.62, 7.08, 4.0, 0.2, 8, "#7A817E")
    add_text(slide, f"{page:02d}", 12.18, 7.06, 0.5, 0.2, 9, "#7A817E", True, "right")


def set_bg(slide, color="#F6F3EE"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def make_ppt(row: dict[str, str]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)
    blank = prs.slide_layouts[6]
    brand = row["브랜드명"]

    # 1. Cover
    s = prs.slides.add_slide(blank)
    set_bg(s)
    s.shapes.add_picture(str(IMAGE_OUT), Inches(6.72), Inches(0.0), width=Inches(6.61), height=Inches(3.72))
    add_text(s, "브랜드별 제안서 샘플", 0.62, 0.72, 3.8, 0.3, 14, "#1F6F5B", True)
    add_text(s, f"{brand}\n영업상품 추천안", 0.62, 1.35, 5.2, 1.28, 42, "#18201E", True)
    add_text(s, "리뷰, FAQ, 장점, 가격 근거를 이용해 어떤 영업상품으로 제안할지 정리한 예시입니다.", 0.66, 3.05, 5.4, 0.64, 18, "#68736F")
    add_badge(s, row["대그룹"], 0.66, 4.18, "#1F6F5B", 2.15)
    add_badge(s, row["세부그룹"], 2.95, 4.18, "#E47C64", 1.75)
    add_text(s, "추천 1순위", 0.66, 5.08, 1.1, 0.2, 11, "#C59043", True)
    add_text(s, "체험단/숏폼 리뷰 콘텐츠", 0.66, 5.38, 4.4, 0.42, 28, "#202624", True)
    add_footer(s, 1)

    # 2. Diagnosis
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FBFAF6")
    add_text(s, "궁중비책은 장벽·진정 메시지가 가장 먼저 보입니다", 0.62, 0.55, 11.5, 0.48, 34, "#18201E", True)
    add_text(s, "브랜드 텍스트에서 반복된 키워드를 제품군, 구매고민, 제안상품으로 연결했습니다.", 0.66, 1.18, 9.8, 0.32, 16, "#68736F")
    cards = [
        ("대그룹", row["대그룹"], row["대그룹설명"]),
        ("세부그룹", row["세부그룹"], "후기와 만족 신호를 콘텐츠 증거로 전환하기 좋은 유형"),
        ("근거 키워드", row["그룹키워드근거"], "장벽·진정·시카·병풀을 중심 메시지로 사용"),
    ]
    for i, (head, title, body) in enumerate(cards):
        x = 0.72 + i * 4.15
        add_rect(s, x, 2.05, 3.55, 2.65, fill="#FFFFFF", line="#DCE3DE")
        add_text(s, head, x + 0.28, 2.33, 1.4, 0.22, 12, "#1F6F5B", True)
        add_text(s, title, x + 0.28, 2.82, 2.85, 0.52, 23, "#202624", True)
        add_text(s, body, x + 0.28, 3.62, 2.85, 0.64, 15, "#68736F")
    add_rect(s, 0.82, 5.45, 11.6, 0.72, fill="#E7F2EC", line="#CFE2D8")
    add_text(s, row["메일첫문장제안"], 1.05, 5.65, 11.1, 0.24, 16, "#1F6F5B", True, "center")
    add_footer(s, 2)

    # 3. Evidence
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_text(s, "근거는 효능 단정보다 구매 전 확인 포인트로 정리합니다", 0.62, 0.55, 11.3, 0.48, 33, "#18201E", True)
    evidence = [
        ("장점", "진정/피부장벽, 가격/혜택, 만족/재구매 신호가 확인됨"),
        ("FAQ", row["기존FAQ근거"]),
        ("후기", row["기존후기근거"]),
        ("경쟁사", row["경쟁사근거"]),
    ]
    for i, (head, body) in enumerate(evidence):
        y = 1.62 + i * 1.18
        add_badge(s, head, 0.78, y + 0.1, "#1F6F5B", 1.05)
        add_text(s, body or "수집 근거 부족", 2.05, y, 9.8, 0.64, 18, "#202624", False, "left", MSO_ANCHOR.MIDDLE)
    add_text(s, "제안서에서는 단점이 비어 있더라도 FAQ와 고객센터 정보를 구매 불안 해소 장표로 배치합니다.", 0.82, 6.36, 11.1, 0.28, 15, "#68736F", True, "center")
    add_footer(s, 3)

    # 4. Recommended sales products
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FBFAF6")
    add_text(s, "이 브랜드는 콘텐츠 증거와 상세페이지 개선을 함께 제안하는 편이 좋습니다", 0.62, 0.55, 11.4, 0.5, 32, "#18201E", True)
    products = [
        ("1순위", "체험단/숏폼 리뷰 콘텐츠", "후기 강점형이라 실제 사용 장면과 만족 신호를 짧게 보여주는 콘텐츠가 가장 자연스럽습니다."),
        ("2순위", "리뷰/FAQ 기반 상세페이지 개선", "FAQ와 장점 근거를 상세페이지의 구매 전 확인 포인트로 재배치합니다."),
        ("3순위", "제안메일용 PPT/이미지 카드 제작", "영업메일 첫 화면에서 그룹 진단과 추천 이유를 한눈에 보여줍니다."),
    ]
    for i, (rank, title, body) in enumerate(products):
        y = 1.62 + i * 1.46
        add_rect(s, 0.72, y, 11.85, 1.03, fill="#FFFFFF", line="#DCE3DE")
        add_badge(s, rank, 1.02, y + 0.33, "#E47C64", 1.0)
        add_text(s, title, 2.25, y + 0.2, 3.6, 0.28, 22, "#202624", True)
        add_text(s, body, 6.0, y + 0.18, 5.75, 0.42, 16, "#68736F")
    add_footer(s, 4)

    # 5. Image direction
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_text(s, "메일 이미지는 세 가지 배지만 보여줘도 충분합니다", 0.62, 0.55, 11.0, 0.48, 34, "#18201E", True)
    s.shapes.add_picture(str(IMAGE_OUT), Inches(0.74), Inches(1.45), width=Inches(6.25), height=Inches(3.52))
    add_text(s, "이미지 카피 배지", 7.35, 1.55, 2.4, 0.3, 15, "#1F6F5B", True)
    for i, label in enumerate(["진정과 장벽 케어", "후기로 확인한 만족도", "구매 혜택 강조", "은은한 사용감"]):
        add_rect(s, 7.35, 2.12 + i * 0.62, 4.1, 0.42, fill="#FFFFFF", line="#DCE3DE")
        add_text(s, label, 7.58, 2.23 + i * 0.62, 3.6, 0.14, 16, "#202624", True)
    add_text(s, "치료/완치처럼 보이는 표현은 피하고, 사용감·성분·후기 근거 중심으로 안전하게 설계합니다.", 7.35, 5.12, 4.28, 0.62, 16, "#68736F")
    add_footer(s, 5)

    # 6. Next step
    s = prs.slides.add_slide(blank)
    set_bg(s, "#18201E")
    add_text(s, "다음 작업은 브랜드별로 같은 구조를 반복 생산하는 것입니다", 0.72, 0.88, 10.6, 0.82, 38, "#FFFFFF", True)
    steps = [
        "검토필요 브랜드 90개는 URL/대표상품을 먼저 정리",
        "근거강도 3 이상 브랜드부터 제안메일 우선 발송",
        "그룹별 이미지 배지를 고정하고 브랜드명만 교체",
        "PPT는 표지, 진단, 근거, 추천상품, 실행안 5장으로 축소 가능",
    ]
    for i, step in enumerate(steps):
        add_text(s, f"{i + 1}", 1.0, 2.25 + i * 0.86, 0.32, 0.26, 17, "#E7C46B", True, "center")
        add_text(s, step, 1.55, 2.2 + i * 0.86, 9.2, 0.32, 21, "#F6F3EE")
    add_text(s, "샘플 파일 기준: ttamtti2_cosmetic_sales_grouping.csv", 0.82, 6.65, 7.2, 0.22, 12, "#B7C2BD")
    add_footer(s, 6)

    prs.save(PPT_OUT)


def main() -> None:
    row = read_brand()
    make_email_image(row)
    make_ppt(row)
    print(f"image={IMAGE_OUT}")
    print(f"ppt={PPT_OUT}")
    print(f"image_bytes={IMAGE_OUT.stat().st_size}")
    print(f"ppt_bytes={PPT_OUT.stat().st_size}")


if __name__ == "__main__":
    main()
