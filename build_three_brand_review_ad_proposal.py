from __future__ import annotations

import csv
import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DATA_DIR = ROOT / "data" / "brand_json_summaries"
LOCAL_CSV = DATA_DIR / "ttamtti2_local_review_analysis.csv"
BRIEF_CSV = DATA_DIR / "ttamtti2_sales_automation_briefs.csv"
OUT_DIR = ROOT / "outputs" / "cosmetic_example"
OUT_PPTX = OUT_DIR / "three_brand_review_ad_proposal.pptx"

BRANDS = ["1.618", "HYCL", "헤이네이처"]

BLACK = RGBColor(16, 16, 16)
MUTED = RGBColor(86, 86, 86)
LIGHT = RGBColor(244, 244, 244)
MID = RGBColor(214, 218, 224)
ORANGE = RGBColor(255, 107, 53)
TEAL = RGBColor(20, 126, 116)
BLUE = RGBColor(42, 91, 215)
WHITE = RGBColor(255, 255, 255)
PALE_ORANGE = RGBColor(255, 244, 239)
PALE_TEAL = RGBColor(235, 247, 245)


BAD_REVIEW_TERMS = {
    "ABOUT",
    "EVENT",
    "BEST",
    "SHOP",
    "COMMUNITY",
    "Login",
    "Join",
    "Cart",
    "자세히보기",
    "검색 가입",
    "브랜드 스토리",
    "Social Contribution",
    "Clean Beauty",
    "글로우픽 앱에서 보기",
}


def read_csv_map(path: Path) -> dict[str, dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return {row["브랜드명"]: row for row in csv.DictReader(f)}


def add_text(slide, text, x, y, w, h, size=18, color=BLACK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def rect(slide, x, y, w, h, fill=LIGHT, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.58, 0.42, 11.6, 0.62, 28, BLACK, True)
    rect(slide, 0.61, 1.23, 11.75, 0.015, MID)
    if subtitle:
        add_text(slide, subtitle, 0.61, 1.35, 11.2, 0.32, 12, MUTED)


def add_footer(slide, page):
    add_text(slide, "브랜드 리뷰 제안서", 0.61, 7.08, 3.0, 0.24, 8, MUTED)
    add_text(slide, f"{page:02d}", 12.0, 7.08, 0.5, 0.24, 8, MUTED, align=PP_ALIGN.RIGHT)


def split_slash(value: str, limit: int | None = None) -> list[str]:
    items = [" ".join(v.strip().split()) for v in (value or "").split(" / ") if v.strip()]
    return items[:limit] if limit else items


def split_comma(value: str, limit: int | None = None) -> list[str]:
    items = [" ".join(v.strip().split()) for v in (value or "").split(",") if v.strip()]
    return items[:limit] if limit else items


def parse_aspect_sentiment(value: str, limit: int = 3) -> list[dict[str, str]]:
    rows = []
    for chunk in split_slash(value):
        match = re.match(r"(.+?):긍정(\d+)%,부정(\d+)%,중립(\d+)%", chunk)
        if not match:
            continue
        aspect, pos, neg, neutral = match.groups()
        pos_i = int(pos)
        neg_i = int(neg)
        if pos_i >= 55 and neg_i < 25:
            use = "메인 소구"
        elif neg_i >= 25:
            use = "보완 메시지"
        else:
            use = "서브 소구"
        rows.append({"aspect": aspect, "pos": pos, "neg": neg, "neutral": neutral, "use": use})
        if len(rows) >= limit:
            break
    return rows


def clean_text(text: str) -> str:
    text = re.sub(r"\([^)]*등록된 네이버 페이 구매평\)", " ", text or "")
    text = re.sub(r"\d{4}[-.]\d{1,2}[-.]\d{1,2}(?:\s+\d{1,2}:\d{2}:\d{2})?", " ", text)
    text = re.sub(r"\b\d{1,2}:\d{2}:\d{2}\b", " ", text)
    text = re.sub(r"[A-Za-z가-힣]\*{2,}", " ", text)
    text = re.sub(r"\b\d{1,3}(?:,\d{3})*원\b", " ", text)
    text = text.replace("더보기", " ").replace("[수정됨]", " ")
    return " ".join(text.split())


def sentence_parts(value: str) -> list[str]:
    parts = []
    for chunk in split_slash(value):
        chunk = clean_text(chunk)
        parts.extend(re.split(r"(?<=[.!?。！？])\s+|(?<=요)\s+|(?<=다)\s+", chunk))
    return [" ".join(part.split()) for part in parts if part.strip()]


def review_score(text: str) -> int:
    score = 0
    good_terms = [
        "좋",
        "촉촉",
        "진정",
        "순하",
        "자극",
        "부드럽",
        "재구매",
        "만족",
        "세안",
        "발림",
        "피부",
        "여드",
        "수분",
        "장벽",
        "흡수",
        "향",
        "선물",
    ]
    score += sum(2 for term in good_terms if term in text)
    if 18 <= len(text) <= 88:
        score += 8
    if len(text) > 112:
        score -= 8
    if any(term in text for term in BAD_REVIEW_TERMS):
        score -= 15
    if text.count(" ") < 2:
        score -= 4
    return score


def best_reviews(local: dict[str, str], limit: int = 3) -> list[str]:
    candidates = []
    for field in ["긍정예시", "대표문장", "부정예시"]:
        candidates.extend(sentence_parts(local.get(field, "")))
    unique = []
    seen = set()
    for item in candidates:
        if len(item) < 12 or item in seen:
            continue
        seen.add(item)
        unique.append(item)
    unique.sort(key=review_score, reverse=True)
    picked = []
    for item in unique:
        if review_score(item) < 2 or len(item) > 92:
            continue
        picked.append(item)
        if len(picked) >= limit:
            break
    return picked or [clean_text(p)[:86] for p in sentence_parts(local.get("대표문장", ""))[:limit]]


def appeal_points(local: dict[str, str], brief: dict[str, str]) -> list[str]:
    text = " ".join(
        [
            local.get("그룹근거키워드", ""),
            local.get("상위키워드", ""),
            local.get("긍정키워드", ""),
            brief.get("공유키워드", ""),
            local.get("대표문장", ""),
        ]
    )
    rules = [
        ("촉촉한 보습감", ["촉촉", "보습", "수분", "건조"]),
        ("민감 피부 진정", ["진정", "어성초", "시카", "자극", "장벽", "여드"]),
        ("산뜻한 사용감", ["산뜻", "흡수", "발림", "끈적", "가벼운"]),
        ("재구매 만족", ["재구매", "만족", "좋아요", "추천", "강추"]),
        ("가격과 혜택", ["가격", "할인", "쿠폰", "혜택", "가성비"]),
        ("클렌징 사용감", ["클렌저", "폼클렌징", "세안", "거품", "모공"]),
        ("톤과 피부결", ["톤", "결", "글로우", "화잘먹", "광"]),
        ("선물성 패키지", ["패키지", "선물", "감성", "고급"]),
    ]
    picked = [label for label, terms in rules if any(term in text for term in terms)]
    return picked[:4] or split_comma(brief.get("공유키워드", ""), 4)


def ad_copy(local: dict[str, str], brief: dict[str, str]) -> list[str]:
    brand = local["브랜드명"]
    points = appeal_points(local, brief)
    group = local.get("대표그룹", "")
    if "진정" in " ".join(points) or "장벽" in group:
        return [
            f"{brand}, 민감한 피부가 먼저 알아보는 진정 케어",
            "자극 걱정은 줄이고 촉촉한 편안함은 오래",
        ]
    if "클렌징" in " ".join(points):
        return [
            f"{brand}, 씻고 난 뒤까지 편안한 클렌징",
            "부드러운 세안감으로 매일 쓰기 좋은 루틴",
        ]
    if "재구매" in " ".join(points):
        return [
            f"{brand}, 다시 찾는 이유가 있는 스킨케어",
            "만족 리뷰가 말해주는 데일리 케어",
        ]
    return [
        f"{brand}, 리뷰가 말하는 핵심 케어",
        f"{points[0] if points else '사용감'}을 중심으로 보여주는 브랜드 경험",
    ]


def image_direction(local: dict[str, str], brief: dict[str, str]) -> list[str]:
    points = " ".join(appeal_points(local, brief))
    group = local.get("대표그룹", "")
    if "진정" in points or "장벽" in group:
        return ["밝은 욕실 선반", "어성초나 초록 식물 포인트", "촉촉한 제형 클로즈업"]
    if "클렌징" in points:
        return ["거품 제형 클로즈업", "세안 후 맑은 피부 연출", "화이트와 민트톤 배경"]
    if "선물" in points:
        return ["패키지 단독 컷", "선물 상자와 자연광", "프리미엄 베이지톤 배경"]
    return ["제품 단독 컷", "은은한 자연광", "피부결이 보이는 사용 장면"]


def segment_rows(local: dict[str, str], brief: dict[str, str]) -> list[tuple[str, str, str]]:
    points = " ".join(appeal_points(local, brief))
    group = local.get("대표그룹", "")
    brand = local["브랜드명"]
    rows: list[tuple[str, str, str]] = []
    if "진정" in points or "장벽" in group:
        rows.append(("민감성 피부 관심 고객", "성분 안정성과 자극 완화", "진정 라인 체험 키트"))
    if "재구매" in points:
        rows.append(("재구매 가능 고객", "만족 리뷰와 반복 구매 이유", "세트 할인 또는 리마인드 쿠폰"))
    if "가격" in points:
        rows.append(("혜택 반응 고객", "가격 부담을 낮추는 구성", "첫 구매 쿠폰과 무료배송"))
    if "클렌징" in points:
        rows.append(("사용감 비교 고객", "세안 후 편안함과 부드러운 제형", "사용감 숏폼 리뷰"))
    if "선물" in points:
        rows.append(("선물 구매 고객", "패키지와 감성 메시지", "기프트 세트 제안"))
    if "톤" in points:
        rows.append(("피부결 개선 관심 고객", "화잘먹과 결 정돈", "전후 비교 이미지"))
    if not rows:
        rows.append((f"{brand} 관심 고객", "리뷰에서 반복된 핵심 장점", "리뷰 기반 상세페이지 개선"))
    return rows[:3]


def tag(slide, text, x, y, w, color):
    rect(slide, x, y, w, 0.34, color)
    add_text(slide, text, x + 0.08, y + 0.07, w - 0.16, 0.15, 8.6, WHITE, True, PP_ALIGN.CENTER)


def metric(slide, label, value, x, y, w, color=BLACK):
    rect(slide, x, y, w, 0.64, LIGHT)
    add_text(slide, value, x + 0.12, y + 0.09, w - 0.24, 0.22, 15, color, True)
    add_text(slide, label, x + 0.12, y + 0.37, w - 0.24, 0.15, 8, MUTED)


def bullet_list(slide, items, x, y, w, h, size=11.3):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(5)
    return box


def small_chip(slide, text, x, y, w):
    rect(slide, x, y, w, 0.28, PALE_TEAL, TEAL)
    add_text(slide, text, x + 0.05, y + 0.06, w - 0.1, 0.12, 8, TEAL, True, PP_ALIGN.CENTER)


def add_brand_slide(prs: Presentation, local: dict[str, str], brief: dict[str, str], page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    brand = local["브랜드명"]
    add_title(slide, f"{brand} geo alt 값", "리뷰 근거, 소구점, 광고 카피, 이미지 방향을 한 장에 정리했습니다.")

    tag(slide, local["대표그룹"], 0.7, 1.8, 2.35, TEAL)
    tag(slide, local["세부유형"], 3.22, 1.8, 1.95, ORANGE)
    tag(slide, f"근거 {local['근거등급']}등급", 5.35, 1.8, 1.22, BLUE)

    metric(slide, "리뷰 문장", local["정제문장수"], 0.72, 2.28, 1.45, BLACK)
    metric(slide, "긍정", local["긍정문장수"], 2.34, 2.28, 1.18, TEAL)
    metric(slide, "주의", local["부정문장수"], 3.68, 2.28, 1.18, ORANGE)
    metric(slide, "우선순위", brief["우선순위점수"], 5.02, 2.28, 1.3, BLUE)

    add_text(slide, "리뷰 근거", 0.72, 3.18, 1.55, 0.25, 15.5, BLACK, True)
    bullet_list(slide, best_reviews(local, 3), 0.78, 3.52, 5.55, 1.08, 10.8)

    add_text(slide, "핵심 소구점", 0.72, 4.95, 1.65, 0.24, 14.5, BLACK, True)
    for i, point in enumerate(appeal_points(local, brief)[:4]):
        small_chip(slide, point, 0.78 + (i % 2) * 2.55, 5.3 + (i // 2) * 0.38, 2.15)

    rect(slide, 6.72, 2.24, 5.0, 1.72, PALE_ORANGE)
    add_text(slide, "광고 카피", 6.98, 2.5, 1.2, 0.2, 15.5, ORANGE, True)
    bullet_list(slide, ad_copy(local, brief), 7.0, 2.92, 4.12, 0.62, 14.0)

    rect(slide, 6.72, 4.38, 5.0, 1.64, PALE_TEAL)
    add_text(slide, "이미지 방향", 6.98, 4.64, 1.35, 0.2, 15.5, TEAL, True)
    bullet_list(slide, image_direction(local, brief), 7.0, 5.04, 4.12, 0.58, 13.0)

    rect(slide, 0.78, 6.2, 5.55, 0.42, PALE_ORANGE)
    add_text(slide, "구매장벽", 0.98, 6.32, 0.95, 0.13, 8.8, ORANGE, True)
    add_text(slide, " / ".join(split_slash(brief["구매장벽"], 2)), 1.94, 6.27, 4.0, 0.18, 9.5, BLACK, True)

    add_footer(slide, page)


def add_segment_strategy_slide(prs: Presentation, locals_by_brand: dict[str, dict[str, str]], briefs_by_brand: dict[str, dict[str, str]], page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "고객 관심사별로 광고 메시지를 나눕니다", "리뷰에서 확인한 소구점을 고객 세그먼트와 프로모션으로 연결합니다.")

    headers = ["브랜드", "주요 고객 세그먼트", "광고 메시지", "프로모션"]
    xs = [0.72, 2.55, 5.35, 8.55]
    ws = [1.45, 2.35, 2.75, 2.95]
    for x, w, header in zip(xs, ws, headers):
        rect(slide, x, 1.95, w, 0.4, BLACK)
        add_text(slide, header, x + 0.08, 2.04, w - 0.16, 0.14, 8.8, WHITE, True, PP_ALIGN.CENTER)

    rows = []
    for brand in BRANDS:
        local = locals_by_brand[brand]
        brief = briefs_by_brand[brand]
        first = segment_rows(local, brief)[0]
        rows.append((brand, *first))

    for idx, row in enumerate(rows):
        y = 2.65 + idx * 1.08
        fill = WHITE if idx % 2 == 0 else LIGHT
        for x, w, text in zip(xs, ws, row):
            rect(slide, x, y, w, 0.72, fill, MID)
            add_text(slide, text, x + 0.1, y + 0.15, w - 0.2, 0.22, 10.3, BLACK, True if x == xs[0] else False)

    rect(slide, 0.82, 6.05, 11.25, 0.48, PALE_ORANGE)
    add_text(slide, "제안 포인트", 1.04, 6.2, 1.2, 0.14, 9.4, ORANGE, True)
    add_text(slide, "같은 리뷰 데이터라도 고객 관심사에 따라 광고 문구와 혜택을 다르게 제안합니다.", 2.18, 6.14, 8.8, 0.2, 12.5, BLACK, True)
    add_footer(slide, page)


def add_segment_detail_slide(prs: Presentation, locals_by_brand: dict[str, dict[str, str]], briefs_by_brand: dict[str, dict[str, str]], page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "세그먼트별 소재 방향까지 바로 제안합니다", "PPT, 이미지 카드, 제안메일에 들어갈 문장을 같은 기준으로 뽑습니다.")

    for idx, brand in enumerate(BRANDS):
        local = locals_by_brand[brand]
        brief = briefs_by_brand[brand]
        x = 0.72 + idx * 4.02
        rect(slide, x, 1.95, 3.55, 4.55, LIGHT)
        add_text(slide, brand, x + 0.22, 2.18, 2.6, 0.28, 17, BLACK, True)
        points = appeal_points(local, brief)[:3]
        add_text(slide, "소구점", x + 0.24, 2.74, 0.8, 0.16, 9.3, TEAL, True)
        bullet_list(slide, points, x + 0.28, 3.03, 3.0, 0.65, 10.2)
        add_text(slide, "세그먼트", x + 0.24, 3.85, 0.85, 0.16, 9.3, ORANGE, True)
        segs = [row[0] for row in segment_rows(local, brief)[:2]]
        bullet_list(slide, segs, x + 0.28, 4.14, 3.0, 0.58, 10.2)
        add_text(slide, "소재 방향", x + 0.24, 4.92, 0.9, 0.16, 9.3, BLUE, True)
        bullet_list(slide, image_direction(local, brief)[:2], x + 0.28, 5.2, 3.0, 0.58, 10.2)

    add_footer(slide, page)


def add_brand_segment_slide(prs: Presentation, local: dict[str, str], brief: dict[str, str], page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    brand = local["브랜드명"]
    add_title(slide, f"{brand} 고객 세그먼트", "리뷰에서 확인한 관심사를 광고 메시지와 프로모션으로 나눴습니다.")

    rows = segment_rows(local, brief)[:3]
    headers = ["고객 세그먼트", "광고 메시지", "프로모션"]
    xs = [0.78, 4.18, 7.72]
    ws = [2.9, 3.05, 3.35]
    for x, w, header in zip(xs, ws, headers):
        rect(slide, x, 1.9, w, 0.48, BLACK)
        add_text(slide, header, x + 0.1, 2.03, w - 0.2, 0.14, 10.2, WHITE, True, PP_ALIGN.CENTER)

    for idx, row in enumerate(rows):
        y = 2.72 + idx * 0.96
        fill = WHITE if idx % 2 == 0 else LIGHT
        for x, w, text in zip(xs, ws, row):
            rect(slide, x, y, w, 0.72, fill, MID)
            add_text(slide, text, x + 0.16, y + 0.22, w - 0.32, 0.2, 11.2, BLACK, True if x == xs[0] else False)

    rect(slide, 0.82, 5.8, 3.65, 0.62, PALE_TEAL)
    add_text(slide, "소재 방향", 1.02, 5.99, 1.0, 0.14, 9.6, TEAL, True)
    add_text(slide, " / ".join(image_direction(local, brief)[:2]), 2.02, 5.93, 2.0, 0.2, 10.8, BLACK, True)

    rect(slide, 4.78, 5.8, 3.35, 0.62, PALE_ORANGE)
    add_text(slide, "구매장벽", 4.98, 5.99, 0.9, 0.14, 9.6, ORANGE, True)
    add_text(slide, " / ".join(split_slash(brief["구매장벽"], 2)), 5.88, 5.91, 1.75, 0.22, 9.8, BLACK, True)

    rect(slide, 8.42, 5.8, 3.2, 0.62, LIGHT)
    add_text(slide, "추천상품", 8.62, 5.99, 0.9, 0.14, 9.6, BLUE, True)
    add_text(slide, split_slash(brief["추천영업상품"], 1)[0], 9.52, 5.91, 1.62, 0.22, 9.7, BLACK, True)

    add_footer(slide, page)


def add_brand_aspect_slide(prs: Presentation, local: dict[str, str], brief: dict[str, str], page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    brand = local["브랜드명"]
    add_title(slide, f"{brand} 항목별 감성 근거", "리뷰 항목별 긍정/부정 반응을 나눠 광고 메시지의 우선순위를 정했습니다.")

    aspect_rows = parse_aspect_sentiment(local.get("항목별감성분포", ""), 5)
    if not aspect_rows:
        aspect_rows = [{"aspect": "리뷰 강점", "pos": "0", "neg": "0", "neutral": "0", "use": "검토 필요"}]

    header_y = 1.98
    table_xs = [0.82, 4.0, 5.35, 6.7]
    table_ws = [2.65, 0.95, 0.95, 1.65]
    for x, w, header in zip(table_xs, table_ws, ["항목", "긍정", "부정", "활용방향"]):
        rect(slide, x, header_y, w, 0.45, BLACK)
        add_text(slide, header, x + 0.08, header_y + 0.13, w - 0.16, 0.13, 9.2, WHITE, True, PP_ALIGN.CENTER)

    for idx, item in enumerate(aspect_rows):
        y = 2.72 + idx * 0.56
        fill = WHITE if idx % 2 == 0 else LIGHT
        values = [item["aspect"], f"{item['pos']}%", f"{item['neg']}%", item["use"]]
        for x, w, value in zip(table_xs, table_ws, values):
            rect(slide, x, y, w, 0.42, fill, MID)
            color = TEAL if value == "메인 소구" else ORANGE if value == "보완 메시지" else BLACK
            add_text(slide, value, x + 0.08, y + 0.12, w - 0.16, 0.13, 9.0, color, True if x in {table_xs[0], table_xs[3]} else False, PP_ALIGN.CENTER if x != table_xs[0] else None)

    rect(slide, 8.72, 2.0, 3.05, 2.35, PALE_ORANGE)
    add_text(slide, "항목별 해석", 8.98, 2.28, 1.1, 0.18, 11.0, ORANGE, True)
    add_text(slide, local.get("ABSC제안근거", "항목별 긍정/부정 반응을 기준으로 광고 메시지를 나눕니다."), 8.98, 2.78, 2.25, 0.86, 11.0, BLACK, True)

    rect(slide, 8.72, 4.75, 3.05, 1.15, PALE_TEAL)
    add_text(slide, "광고 적용", 8.98, 5.02, 1.0, 0.16, 10.2, TEAL, True)
    add_text(slide, "메인 소구는 전면 카피로, 보완 항목은 FAQ/혜택/상세페이지 문구로 분리합니다.", 8.98, 5.42, 2.25, 0.34, 10.2, BLACK, True)

    add_text(slide, "제안 문장", 0.82, 6.47, 0.95, 0.14, 9.3, ORANGE, True)
    add_text(
        slide,
        f"{brand}는 긍정 비율이 높은 항목을 메인 소구로 쓰고, 부정 또는 호불호 항목은 보완 메시지로 분리하는 제안이 적합합니다.",
        1.82,
        6.4,
        9.8,
        0.22,
        11.2,
        BLACK,
        True,
    )
    add_footer(slide, page)


def add_cover(prs: Presentation, locals_by_brand: dict[str, dict[str, str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, "브랜드별 리뷰 기반\n광고 제안 샘플", 0.66, 0.88, 7.4, 1.4, 36, BLACK, True)
    add_text(slide, "리뷰에서 고객 반응을 뽑고, 그 반응을 광고 카피와 이미지 방향으로 연결했습니다.", 0.7, 2.62, 7.4, 0.46, 15, MUTED)
    rect(slide, 8.35, 0.86, 4.1, 5.55, LIGHT)
    add_text(slide, "포함 브랜드", 8.7, 1.18, 2.2, 0.28, 17, BLACK, True)
    for idx, brand in enumerate(BRANDS, start=1):
        row = locals_by_brand[brand]
        y = 1.78 + (idx - 1) * 1.15
        add_text(slide, f"0{idx}", 8.72, y, 0.55, 0.34, 18, ORANGE, True)
        add_text(slide, row["브랜드명"], 9.35, y, 2.5, 0.34, 18, BLACK, True)
        add_text(slide, row["대표그룹"], 9.35, y + 0.38, 2.7, 0.22, 10, MUTED)
    add_footer(slide, 1)


def add_close(prs: Presentation, locals_by_brand: dict[str, dict[str, str]], briefs_by_brand: dict[str, dict[str, str]], page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "리뷰가 말하는 소구점이 광고 방향을 결정합니다", "리뷰 문장, 소구점, 추천상품을 같이 보면 브랜드별 제안이 더 선명해집니다.")
    for idx, brand in enumerate(BRANDS):
        local = locals_by_brand[brand]
        brief = briefs_by_brand[brand]
        y = 1.92 + idx * 1.18
        rect(slide, 0.82, y, 11.4, 0.86, LIGHT)
        add_text(slide, brand, 1.08, y + 0.2, 1.65, 0.24, 16, BLACK, True)
        add_text(slide, " / ".join(appeal_points(local, brief)[:2]), 3.0, y + 0.2, 3.1, 0.22, 12.2, TEAL, True)
        add_text(slide, split_slash(brief["추천영업상품"], 1)[0], 6.48, y + 0.2, 4.7, 0.22, 12.2, BLACK)
    rect(slide, 0.82, 5.78, 11.4, 0.62, PALE_ORANGE)
    add_text(slide, "다음 작업", 1.08, 5.98, 1.1, 0.18, 10.5, ORANGE, True)
    add_text(slide, "브랜드별 리뷰 근거를 유지한 채 이미지 카드와 제안메일 본문까지 같은 구조로 확장합니다.", 2.18, 5.9, 8.7, 0.28, 14, BLACK, True)
    add_footer(slide, page)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    locals_by_brand = read_csv_map(LOCAL_CSV)
    briefs_by_brand = read_csv_map(BRIEF_CSV)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs, locals_by_brand)
    page = 2
    for brand in BRANDS:
        add_brand_slide(prs, locals_by_brand[brand], briefs_by_brand[brand], page)
        page += 1
        add_brand_segment_slide(prs, locals_by_brand[brand], briefs_by_brand[brand], page)
        page += 1
        add_brand_aspect_slide(prs, locals_by_brand[brand], briefs_by_brand[brand], page)
        page += 1
    add_close(prs, locals_by_brand, briefs_by_brand, page)

    try:
        prs.save(OUT_PPTX)
        print(OUT_PPTX)
    except PermissionError:
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        fallback = OUT_DIR / f"three_brand_review_ad_proposal_absc_{stamp}.pptx"
        prs.save(fallback)
        print(fallback)


if __name__ == "__main__":
    main()
