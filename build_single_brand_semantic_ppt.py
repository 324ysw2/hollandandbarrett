from __future__ import annotations

import argparse
import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
CSV_PATH = ROOT / "data" / "brand_json_summaries" / "ttamtti2_semantic_brand_recommendations.csv"
OUT_DIR = ROOT / "outputs" / "cosmetic_example"

BLACK = RGBColor(14, 14, 14)
MUTED = RGBColor(86, 86, 86)
LIGHT = RGBColor(241, 241, 241)
LINE = RGBColor(190, 194, 200)
ORANGE = RGBColor(255, 107, 53)
TEAL = RGBColor(20, 126, 116)
BLUE = RGBColor(42, 91, 215)
PALE_TEAL = RGBColor(235, 247, 245)
PALE_ORANGE = RGBColor(255, 244, 239)
WHITE = RGBColor(255, 255, 255)


def read_brand(brand_name: str) -> dict[str, str]:
    with CSV_PATH.open("r", encoding="utf-8-sig", newline="") as f:
        for row in csv.DictReader(f):
            if row["브랜드명"] == brand_name:
                return row
    raise SystemExit(f"브랜드를 찾을 수 없습니다: {brand_name}")


def add_text(slide, text, x, y, w, h, size=16, color=BLACK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def rect(slide, x, y, w, h, fill=LIGHT, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def tag(slide, text, x, y, w, color):
    rect(slide, x, y, w, 0.36, color)
    add_text(slide, text, x + 0.08, y + 0.08, w - 0.16, 0.14, 8.5, WHITE, True, PP_ALIGN.CENTER)


def bullet_list(slide, items, x, y, w, h, size=12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
    return box


def split_slash(value: str, limit: int | None = None) -> list[str]:
    items = [" ".join(v.strip().split()) for v in (value or "").split(" / ") if v.strip()]
    return items[:limit] if limit else items


def split_comma(value: str, limit: int = 8) -> list[str]:
    return [" ".join(v.strip().split()) for v in (value or "").split(",") if v.strip()][:limit]


def shorten(text: str, limit: int) -> str:
    text = " ".join((text or "").split())
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "..."


def build_ppt(row: dict[str, str], output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    brand = row["브랜드명"]
    add_text(slide, f"{brand} 의미 기반 영업 제안", 0.62, 0.48, 8.6, 0.58, 30, BLACK, True)
    add_text(slide, "리뷰/FAQ/키워드에서 가까운 브랜드를 찾고, 제안할 영업상품을 자동 추천한 샘플입니다.", 0.66, 1.18, 9.6, 0.28, 12, MUTED)
    rect(slide, 0.66, 1.58, 11.95, 0.015, LINE)

    tag(slide, row["대표그룹"], 0.72, 1.95, 2.35, TEAL)
    tag(slide, row["세부유형"], 3.25, 1.95, 1.8, ORANGE)
    tag(slide, f"근거 {row['근거등급']}등급", 5.22, 1.95, 1.22, BLUE)

    add_text(slide, "공유 키워드", 0.72, 2.72, 1.6, 0.28, 17, BLACK, True)
    for idx, kw in enumerate(split_comma(row["공유키워드"], 10)):
        x = 0.74 + (idx % 5) * 1.1
        y = 3.18 + (idx // 5) * 0.46
        rect(slide, x, y, 0.92, 0.3, PALE_TEAL, TEAL)
        add_text(slide, kw, x + 0.04, y + 0.065, 0.84, 0.12, 8, TEAL, True, PP_ALIGN.CENTER)

    add_text(slide, "유사 브랜드", 6.95, 2.72, 1.8, 0.28, 17, BLACK, True)
    similar = split_slash(row["유사브랜드"], 5)
    sims = split_slash(row["유사도"], 5)
    for idx, name in enumerate(similar):
        y = 3.12 + idx * 0.42
        add_text(slide, f"{idx + 1}. {name}", 7.02, y, 2.2, 0.2, 11.5, BLACK, True if idx == 0 else False)
        if idx < len(sims):
            add_text(slide, sims[idx], 9.45, y, 0.65, 0.2, 10, MUTED, align=PP_ALIGN.RIGHT)

    rect(slide, 0.72, 4.62, 5.82, 1.12, LIGHT)
    add_text(slide, "의미기반추천영업상품", 0.96, 4.82, 2.4, 0.22, 12, ORANGE, True)
    bullet_list(slide, split_slash(row["의미기반추천영업상품"], 4), 0.98, 5.14, 5.2, 0.42, 11.5)

    rect(slide, 6.94, 4.62, 5.22, 1.12, PALE_ORANGE)
    add_text(slide, "추천근거", 7.18, 4.82, 1.2, 0.22, 12, ORANGE, True)
    add_text(slide, shorten(row["추천근거"], 128), 7.18, 5.14, 4.62, 0.36, 11.5, BLACK, True)

    add_text(slide, "제안메일용 한 줄", 0.74, 6.28, 2.0, 0.24, 13, TEAL, True)
    one_liner = (
        f"{brand}는 {split_comma(row['공유키워드'], 3)[0]}, "
        f"{split_comma(row['공유키워드'], 3)[1]}, "
        f"{split_comma(row['공유키워드'], 3)[2]} 신호가 있어 "
        f"{split_slash(row['의미기반추천영업상품'], 1)[0]} 제안이 적합합니다."
    )
    add_text(slide, one_liner, 2.55, 6.26, 8.9, 0.26, 13, BLACK, True)
    add_text(slide, "Local semantic brand recommender | no API", 0.66, 7.1, 3.4, 0.18, 8, MUTED)

    prs.save(output)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--brand", default="헤이네이처")
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()

    row = read_brand(args.brand)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    output = args.output or OUT_DIR / f"{args.brand}_semantic_sales_onepage.pptx"
    build_ppt(row, output)
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
