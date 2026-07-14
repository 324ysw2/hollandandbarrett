from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(r"C:\projects\db\본부장님께\플루")
OUT_PPTX = OUT_DIR / "플루_GEO_SEO_고도화_전략_20260710.pptx"


NAVY = RGBColor(15, 33, 61)
TEAL = RGBColor(0, 132, 128)
MINT = RGBColor(228, 247, 244)
CORAL = RGBColor(239, 102, 86)
INK = RGBColor(27, 31, 40)
MUTED = RGBColor(92, 99, 112)
LINE = RGBColor(216, 221, 230)
LIGHT = RGBColor(247, 249, 252)
WHITE = RGBColor(255, 255, 255)
AMBER = RGBColor(247, 176, 62)
BLUE = RGBColor(59, 111, 216)


def set_font(run, size=14, color=INK, bold=False):
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(slide, text, x, y, w, h, size=14, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align is not None:
        p.alignment = align
    return box


def rect(slide, x, y, w, h, fill=LIGHT, line=None, radius=False):
    shape_type = 5 if radius else 1
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_title(slide, title, subtitle=None, page=None):
    add_text(slide, title, 0.55, 0.38, 11.5, 0.55, 24, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.96, 11.6, 0.34, 11, MUTED)
    rect(slide, 0.58, 1.37, 12.15, 0.015, LINE)
    if page:
        add_text(slide, f"{page:02d}", 12.18, 0.44, 0.45, 0.24, 8, MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide):
    add_text(slide, "PLU GEO/SEO strategy | Review-driven structured content", 0.58, 7.12, 5.0, 0.22, 8, MUTED)


def bullet_list(slide, items, x, y, w, h, size=12, color=INK, gap=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
    return box


def pill(slide, text, x, y, w, color=TEAL):
    rect(slide, x, y, w, 0.32, color, color, radius=True)
    add_text(slide, text, x + 0.12, y + 0.065, w - 0.24, 0.15, 8.5, WHITE, True, PP_ALIGN.CENTER)


def card(slide, title, body, x, y, w, h, accent=TEAL):
    rect(slide, x, y, w, h, WHITE, LINE)
    rect(slide, x, y, 0.08, h, accent, accent)
    add_text(slide, title, x + 0.22, y + 0.18, w - 0.36, 0.28, 13, NAVY, True)
    bullet_list(slide, body, x + 0.22, y + 0.58, w - 0.38, h - 0.75, 10.5, INK, 3)


def metric_card(slide, title, value, note, x, y, w, color):
    rect(slide, x, y, w, 1.02, WHITE, LINE)
    add_text(slide, title, x + 0.16, y + 0.12, w - 0.3, 0.2, 9.5, MUTED, True)
    add_text(slide, value, x + 0.16, y + 0.34, 0.76, 0.36, 23, color, True)
    add_text(slide, note, x + 0.95, y + 0.36, w - 1.1, 0.32, 9.5, INK)


def code_box(slide, lines, x, y, w, h):
    rect(slide, x, y, w, h, RGBColor(13, 43, 60), RGBColor(13, 43, 60))
    text = "\n".join(lines)
    add_text(slide, text, x + 0.22, y + 0.22, w - 0.44, h - 0.44, 13, WHITE)


def table_header(slide, headers, xs, y, hs, color=NAVY):
    for header, x, w in zip(headers, xs, hs):
        rect(slide, x, y, w, 0.36, color, color)
        add_text(slide, header, x + 0.07, y + 0.08, w - 0.14, 0.14, 8.5, WHITE, True, PP_ALIGN.CENTER)


def table_row(slide, values, xs, y, ws, fill=WHITE, bold_first=True):
    for idx, (value, x, w) in enumerate(zip(values, xs, ws)):
        rect(slide, x, y, w, 0.5, fill, LINE)
        add_text(slide, value, x + 0.08, y + 0.08, w - 0.16, 0.22, 8.8, INK, bold_first and idx == 0)


def build_ppt() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, 13.333, 7.5, WHITE)
    rect(slide, 0, 0, 13.333, 1.0, NAVY, NAVY)
    add_text(slide, "AI 검색 시대를 대비한", 0.7, 1.35, 5.8, 0.38, 20, TEAL, True)
    add_text(slide, "플루 GEO/SEO\n고도화 전략", 0.68, 1.78, 6.9, 1.42, 38, NAVY, True)
    add_text(slide, "리뷰와 구매동기를 구조화해 AI가 추천하는 바디케어 브랜드로 만드는 방법", 0.72, 3.42, 6.9, 0.55, 16, MUTED)
    rect(slide, 8.05, 1.42, 4.3, 4.28, MINT, MINT)
    add_text(slide, "핵심 전환", 8.45, 1.82, 2.4, 0.32, 16, NAVY, True)
    bullet_list(slide, [
        "키워드 삽입 중심 SEO에서 벗어나기",
        "제품, 리뷰, 사용맥락을 구조화하기",
        "AI가 답변에 인용할 수 있는 근거 만들기",
        "저자극, 향, 두피, 루틴 확장 키워드 강화"
    ], 8.48, 2.34, 3.35, 2.25, 13)
    pill(slide, "Product", 8.5, 5.02, 0.95, TEAL)
    pill(slide, "Review", 9.6, 5.02, 0.95, BLUE)
    pill(slide, "FAQ", 10.7, 5.02, 0.8, CORAL)
    pill(slide, "HowTo", 11.65, 5.02, 0.8, AMBER)
    add_text(slide, "2026.07.10", 0.72, 6.75, 2.3, 0.22, 10, MUTED)
    add_footer(slide)

    # 2. Why GEO now
    slide = prs.slides.add_slide(blank)
    add_title(slide, "왜 지금 GEO가 중요한가", "소비자의 검색은 키워드 입력에서 AI에게 질문하는 방식으로 이동하고 있습니다.", 2)
    card(slide, "기존 검색", [
        "소비자가 직접 검색 결과를 비교",
        "상품명, 가격, 후기, 광고 문구를 하나씩 확인",
        "상세페이지의 가독성과 광고 노출이 중요"
    ], 0.7, 1.85, 3.75, 2.25, BLUE)
    card(slide, "AI 검색", [
        "AI가 먼저 후보 상품을 요약하고 추천",
        "홈페이지의 구조화된 정보와 리뷰 근거를 읽음",
        "질문 의도와 제품 포지션의 연결성이 중요"
    ], 4.85, 1.85, 3.75, 2.25, TEAL)
    card(slide, "플루의 기회", [
        "바디스크럽 대표 인식은 이미 보유",
        "저자극, 향, 두피, 루틴 확장으로 답변 영역 확대",
        "리뷰 기반 FAQ와 스키마 적용 시 AI 인용 가능성 상승"
    ], 9.0, 1.85, 3.6, 2.25, CORAL)
    rect(slide, 0.7, 4.72, 11.9, 1.18, LIGHT, LINE)
    add_text(slide, "소비자가 AI에게 묻는 질문 예시", 0.95, 4.92, 2.35, 0.26, 13, NAVY, True)
    bullet_list(slide, [
        "올리브영에서 살만한 저자극 바디스크럽 추천해줘",
        "두피 냄새랑 각질에 좋은 스크럽 있어?",
        "향 오래가는 바디워시 뭐가 좋아?",
        "구취제거 치약 중에 데일리로 쓰기 좋은 거 알려줘"
    ], 3.25, 4.87, 8.9, 0.7, 11, INK, 1)
    add_footer(slide)

    # 3. Product positioning
    slide = prs.slides.add_slide(blank)
    add_title(slide, "플루 제품 포지션은 '상품명'보다 '고민 순간'으로 묶어야 합니다", "AI 추천은 제품군보다 사용자가 처한 상황과 해결하고 싶은 문제에 반응합니다.", 3)
    card(slide, "오리지널 바디스크럽", [
        "검색어: 플루 바디스크럽, 올리브영 바디스크럽, 화이트머스크 바디스크럽",
        "구매동기: 팔꿈치, 무릎, 등, 다리 각질",
        "상황: 여름 노출 전 피부결 정리, 올영 세일 장바구니"
    ], 0.65, 1.75, 5.95, 1.25, TEAL)
    card(slide, "퍼퓸드 바디스크럽/스크럽 바디워시", [
        "검색어: 향 좋은 바디스크럽, 살냄새 바디스크럽, 퍼퓸 바디워시",
        "구매동기: 향 지속, 샤워 루틴, 기분전환",
        "상황: 바디미스트 전 단계, 선물 구매"
    ], 6.85, 1.75, 5.85, 1.25, CORAL)
    card(slide, "두피 스케일링 스크럽", [
        "검색어: 두피 스케일링 스크럽, 두피 열감, 두피 각질",
        "구매동기: 여름 두피 냄새, 피지, 각질",
        "상황: 샴푸만으로 부족하다는 문제의식"
    ], 0.65, 3.45, 5.95, 1.25, BLUE)
    card(slide, "바디로션/대용량 바디워시/구강케어", [
        "검색어: 플루 바디로션, 대용량 바디워시, 구취제거 치약",
        "구매동기: 기존 스크럽 고객의 루틴 확장",
        "상황: 1+1, 세트 구매, 가족용 데일리 케어"
    ], 6.85, 3.45, 5.85, 1.25, AMBER)
    rect(slide, 0.65, 5.38, 12.05, 0.62, NAVY, NAVY)
    add_text(slide, "정리: 플루는 '각질 제거 브랜드'에서 '샤워, 두피, 구강까지 이어지는 데일리 클린케어 루틴 브랜드'로 확장해 설명해야 합니다.", 0.95, 5.56, 11.45, 0.22, 14, WHITE, True)
    add_footer(slide)

    # 4. Keyword and motive enhancement
    slide = prs.slides.add_slide(blank)
    add_title(slide, "검색어와 구매동기 고도화", "핵심은 검색어를 고민, 구매상황, 필요한 증거까지 연결하는 것입니다.", 4)
    headers = ["검색어 축", "사용자 고민", "구매상황", "페이지에 필요한 증거"]
    xs = [0.65, 3.1, 5.95, 8.9]
    ws = [2.25, 2.65, 2.75, 3.75]
    table_header(slide, headers, xs, 1.75, ws)
    rows = [
        ["저자극 바디스크럽", "스크럽은 따갑지 않을까", "예민한 피부도 각질 관리 필요", "성분, 테스트, 부위별 사용법, 민감 피부 FAQ"],
        ["향 좋은 바디스크럽", "향이 금방 날아가지 않을까", "샤워 후 잔향과 기분전환", "향 노트, 리뷰 표현, 바디미스트 루틴"],
        ["두피 스케일링", "샴푸만으로 개운하지 않음", "여름 냄새, 열감, 각질 관리", "사용 주기, 샴푸와 차이, 두피 타입별 FAQ"],
        ["바디워시 추천", "너무 흔해서 고르기 어려움", "향, 피부결, 대용량, 세일", "향/스크럽/보습 기준 비교표"],
        ["구취제거 치약", "상쾌함이 오래갈까", "출근, 외출, 식후 데일리 루틴", "구취 원인 FAQ, 사용감 리뷰, 루틴 연결"]
    ]
    for i, row in enumerate(rows):
        table_row(slide, row, xs, 2.22 + i * 0.62, ws, LIGHT if i % 2 else WHITE)
    rect(slide, 0.65, 5.65, 12.0, 0.52, MINT, MINT)
    add_text(slide, "GEO 관점: AI가 추천할 수 있는 문장은 '키워드'가 아니라 '누구에게, 어떤 상황에서, 왜 맞는지'가 분명한 문장입니다.", 0.92, 5.79, 11.4, 0.18, 12.5, NAVY, True)
    add_footer(slide)

    # 5. Weak keyword strategy
    slide = prs.slides.add_slide(blank)
    add_title(slide, "약한 키워드별 GEO 해결 전략", "광고연관지수는 있지만 클릭기대지수가 낮은 영역은 랜딩 근거와 AI 친화 콘텐츠가 필요합니다.", 5)
    metric_card(slide, "저자극 바디스크럽", "80/68", "저자극 표현은 있으나 성분, 테스트, 민감 피부 FAQ 근거 보강 필요", 0.7, 1.75, 5.85, TEAL)
    metric_card(slide, "두피 스케일링 스크럽", "78/66", "제품은 있으나 바디스크럽 대표 인식이 강해 카테고리 교육 필요", 6.8, 1.75, 5.85, BLUE)
    metric_card(slide, "바디워시 추천", "58/45", "경쟁 폭이 넓어 향 좋은 바디워시, 스크럽 바디워시로 좁혀야 함", 0.7, 3.15, 5.85, CORAL)
    metric_card(slide, "구취제거 치약", "42/38", "플루와 구강케어의 연결이 약해 데일리 클린 루틴으로 묶어야 함", 6.8, 3.15, 5.85, AMBER)
    card(slide, "해결 원칙", [
        "각 키워드별 전용 랜딩 섹션 생성",
        "리뷰에서 반복되는 표현을 FAQ와 Review schema에 반영",
        "Product, FAQPage, HowTo, BreadcrumbList를 함께 적용",
        "AI 질문형 문장에 바로 답하는 문단을 상세페이지 상단에 배치"
    ], 0.7, 4.75, 11.95, 1.22, NAVY)
    add_footer(slide)

    # 6. Homepage structure
    slide = prs.slides.add_slide(blank)
    add_title(slide, "AI가 이해하기 쉬운 홈페이지 구조", "플루 공식몰은 자체개발형으로 보고, 구조화 데이터를 직접 삽입하는 방식이 적합합니다.", 6)
    code_box(slide, [
        '<h1>플루 바디스크럽 추천 | 저자극 각질케어와 향 좋은 샤워 루틴</h1>',
        '',
        '<img alt="플루 오리지널 바디스크럽 화이트머스크 향 제품 이미지" />',
        '',
        '<script type="application/ld+json">',
        '  Product, Offer, FAQPage, Review, BreadcrumbList, HowTo',
        '</script>'
    ], 0.72, 1.75, 6.1, 3.08)
    card(slide, "코드에 넣어야 할 의미", [
        "상품명, 브랜드, 카테고리, 용량, 가격, 구매 가능 여부",
        "향, 사용부위, 피부/두피 타입, 사용 주기",
        "리뷰 핵심 문장과 반복 언급 키워드",
        "소비자 질문과 답변, 사용법, 비교표"
    ], 7.18, 1.75, 5.18, 1.75, TEAL)
    card(slide, "AI가 읽어야 하는 결론", [
        "플루는 어떤 문제를 해결하는가",
        "누구에게 추천할 수 있는가",
        "리뷰 근거는 무엇인가",
        "다른 바디케어 제품과 무엇이 다른가"
    ], 7.18, 3.78, 5.18, 1.45, CORAL)
    rect(slide, 0.72, 5.55, 11.65, 0.54, NAVY, NAVY)
    add_text(slide, "AI는 예쁜 상세페이지를 보는 것이 아니라, 코드 안에 정리된 상품명, 카테고리, 리뷰, 구매동기, 사용상황, FAQ를 읽고 추천 여부를 판단합니다.", 1.0, 5.7, 11.1, 0.18, 12.2, WHITE, True)
    add_footer(slide)

    # 7. Review data
    slide = prs.slides.add_slide(blank)
    add_title(slide, "리뷰 분석은 별점이 아니라 구매동기 데이터입니다", "리뷰에서 뽑은 언어를 상세페이지와 구조화 데이터의 근거 문장으로 재활용합니다.", 7)
    card(slide, "리뷰에서 추출할 항목", [
        "구매 전 고민: 각질, 냄새, 열감, 향 지속, 구취",
        "사용 부위: 팔꿈치, 무릎, 등, 다리, 두피, 구강",
        "만족 포인트: 부드러움, 개운함, 잔향, 피부결",
        "불만/주의 포인트: 자극감, 사용 주기, 향 호불호"
    ], 0.7, 1.78, 5.8, 1.82, BLUE)
    card(slide, "콘텐츠로 바꾸는 방식", [
        "상세페이지 상단: 대표 구매동기 3개 노출",
        "FAQ: 실제 질문형 검색어로 작성",
        "Review schema: 반복 리뷰 문장 요약",
        "HowTo: 사용순서와 주기 명시"
    ], 6.85, 1.78, 5.75, 1.82, TEAL)
    rect(slide, 0.7, 4.12, 11.9, 1.55, LIGHT, LINE)
    add_text(slide, "예시 변환", 0.98, 4.34, 1.1, 0.24, 13, NAVY, True)
    bullet_list(slide, [
        '"팔꿈치가 부드러워졌다" -> 각질 제거/피부결 개선 근거',
        '"향이 오래간다" -> 퍼퓸 바디워시/향 좋은 바디스크럽 근거',
        '"두피가 개운하다" -> 두피 스케일링/여름 두피 관리 근거',
        '"상쾌함이 남는다" -> 구취제거 치약/데일리 구강 루틴 근거'
    ], 2.05, 4.28, 9.7, 0.88, 11.5, INK, 2)
    add_footer(slide)

    # 8. Page blueprint
    slide = prs.slides.add_slide(blank)
    add_title(slide, "플루 상세페이지 GEO 블루프린트", "하나의 상품페이지가 검색광고, 네이버 검색, 구글, AI 답변에 동시에 대응하도록 설계합니다.", 8)
    headers = ["페이지 영역", "넣어야 할 내용", "적용 스키마"]
    xs = [0.72, 3.45, 8.35]
    ws = [2.45, 4.55, 4.0]
    table_header(slide, headers, xs, 1.75, ws)
    rows = [
        ["상단 H1/H2", "상품명 + 구매동기 + 핵심 검색어", "Product"],
        ["리뷰 근거", "반복 리뷰 표현, 만족/주의 포인트", "Review"],
        ["사용법", "사용 부위, 주기, 순서, 함께 쓰는 제품", "HowTo"],
        ["질문 답변", "AI 질문형 검색어에 대한 답변", "FAQPage"],
        ["카테고리 경로", "바디케어 > 바디스크럽 > 저자극", "BreadcrumbList"],
        ["비교표", "오리지널, 퍼퓸드, 두피, 바디워시 비교", "ItemList"]
    ]
    for i, row in enumerate(rows):
        table_row(slide, row, xs, 2.2 + i * 0.55, ws, LIGHT if i % 2 else WHITE)
    add_footer(slide)

    # 9. Roadmap
    slide = prs.slides.add_slide(blank)
    add_title(slide, "실행 로드맵", "빠르게 고칠 수 있는 영역부터 적용하고, 리뷰 데이터가 쌓일수록 AI 답변 근거를 강화합니다.", 9)
    phases = [
        ("1주차", "정리", ["제품군별 키워드/구매동기 확정", "리뷰 문장 분류", "약한 키워드별 FAQ 초안"]),
        ("2주차", "구조화", ["Product/FAQ/Review schema 설계", "상세페이지 H1/H2/alt 정비", "카테고리 breadcrumb 정리"]),
        ("3주차", "랜딩", ["저자극/두피/향/구취 전용 섹션 구축", "비교표와 사용법 콘텐츠 추가", "검색광고 랜딩 연결"]),
        ("4주차", "측정", ["AI 질문 노출 여부 점검", "검색광고 CTR/CVR 비교", "리뷰 기반 FAQ 업데이트"])
    ]
    for i, (week, title, items) in enumerate(phases):
        x = 0.72 + i * 3.05
        rect(slide, x, 1.85, 2.65, 3.8, WHITE, LINE)
        rect(slide, x, 1.85, 2.65, 0.58, [TEAL, BLUE, CORAL, AMBER][i], [TEAL, BLUE, CORAL, AMBER][i])
        add_text(slide, week, x + 0.18, 2.02, 0.75, 0.16, 10, WHITE, True)
        add_text(slide, title, x + 0.9, 1.98, 1.45, 0.2, 13, WHITE, True)
        bullet_list(slide, items, x + 0.22, 2.68, 2.2, 1.88, 10.7, INK, 4)
    rect(slide, 0.72, 6.0, 11.8, 0.46, MINT, MINT)
    add_text(slide, "목표: 플루가 AI 답변에서 '저자극 바디스크럽', '향 좋은 바디워시', '두피 스케일링', '데일리 클린케어 루틴'의 추천 후보로 읽히게 만드는 것", 0.98, 6.13, 11.25, 0.16, 11.7, NAVY, True)
    add_footer(slide)

    # 10. Final message
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, 13.333, 7.5, NAVY, NAVY)
    add_text(slide, "최종 메시지", 0.85, 0.78, 2.5, 0.36, 20, MINT, True)
    add_text(slide, "플루의 GEO 전략은\n키워드를 많이 넣는 작업이 아닙니다.", 0.82, 1.45, 9.2, 1.25, 34, WHITE, True)
    add_text(slide, "제품별 리뷰와 구매동기를 구조화해 AI가 플루를 정확히 이해하고 추천할 수 있게 만드는 작업입니다.", 0.86, 3.02, 8.4, 0.78, 20, WHITE)
    rect(slide, 0.86, 4.28, 11.5, 1.25, RGBColor(23, 52, 89), RGBColor(23, 52, 89))
    bullet_list(slide, [
        "저자극: 성분/테스트/민감 피부 FAQ로 신뢰 강화",
        "향: 퍼퓸 루틴과 리뷰 표현으로 클릭 기대 상승",
        "두피: 샴푸만으로 부족한 여름 두피 관리 카테고리 교육",
        "구강케어: 바디케어 고객의 데일리 클린 루틴으로 확장"
    ], 1.18, 4.5, 10.7, 0.7, 12.2, WHITE, 1)
    add_text(slide, "PLU | GEO/SEO 고도화 제안", 0.9, 6.72, 3.2, 0.24, 10, MINT)

    prs.save(OUT_PPTX)
    print(OUT_PPTX)


if __name__ == "__main__":
    build_ppt()
