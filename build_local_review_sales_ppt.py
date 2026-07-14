from __future__ import annotations

import json
from collections import OrderedDict
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DATA_DIR = ROOT / "data" / "brand_json_summaries"
SUMMARY_PATH = DATA_DIR / "ttamtti2_local_review_analysis_summary.json"
CSV_PATH = DATA_DIR / "ttamtti2_local_review_analysis.csv"
OUT_DIR = ROOT / "outputs" / "cosmetic_example"
OUT_PPTX = OUT_DIR / "ttamtti2_local_review_sales_sample.pptx"


BLACK = RGBColor(14, 14, 14)
MUTED = RGBColor(86, 86, 86)
LIGHT = RGBColor(239, 239, 239)
LINE = RGBColor(188, 192, 198)
ORANGE = RGBColor(255, 107, 53)
TEAL = RGBColor(20, 126, 116)
BLUE = RGBColor(42, 91, 215)
WHITE = RGBColor(255, 255, 255)


def read_csv_rows() -> list[dict[str, str]]:
    import csv

    with CSV_PATH.open("r", encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def add_text(slide, text, x, y, w, h, size=20, color=BLACK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.42, 11.55, 0.72, 28, BLACK, True)
    slide.shapes.add_shape(1, Inches(0.58), Inches(1.27), Inches(11.85), Inches(0.01)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = LINE
    slide.shapes[-1].line.color.rgb = LINE
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.35, 11.2, 0.34, 12, MUTED)


def add_footer(slide, page):
    add_text(slide, "Local review analysis | no API", 0.58, 7.08, 3.2, 0.25, 8, MUTED)
    add_text(slide, f"{page:02d}", 12.1, 7.08, 0.5, 0.25, 8, MUTED, align=PP_ALIGN.RIGHT)


def rect(slide, x, y, w, h, fill=LIGHT, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def metric(slide, label, value, x, y, w, color=BLACK):
    rect(slide, x, y, w, 1.05, LIGHT)
    add_text(slide, value, x + 0.18, y + 0.15, w - 0.36, 0.38, 24, color, True)
    add_text(slide, label, x + 0.18, y + 0.62, w - 0.36, 0.28, 10, MUTED)


def bar(slide, label, value, max_value, x, y, w, color=BLACK):
    add_text(slide, label, x, y - 0.03, 2.25, 0.25, 10, BLACK)
    add_text(slide, str(value), x + 2.38, y - 0.03, 0.55, 0.25, 10, MUTED, align=PP_ALIGN.RIGHT)
    rect(slide, x + 3.05, y, w, 0.13, RGBColor(230, 230, 230))
    if max_value:
        rect(slide, x + 3.05, y, w * value / max_value, 0.13, color)


def bullet_list(slide, items, x, y, w, h, size=13):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.level = 0
        p.space_after = Pt(8)
    return box


def clean_sentence(text: str, limit: int = 130) -> str:
    text = " ".join((text or "").replace("/", " ").split())
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "..."


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    summary = json.loads(SUMMARY_PATH.read_text(encoding="utf-8"))
    rows = read_csv_rows()
    sample = next((r for r in rows if r["브랜드명"] == "와더스킨"), None)
    if sample is None:
        sample = rows[0]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, 13.333, 7.5, WHITE)
    add_text(slide, "화장품 브랜드 리뷰 기반\n영업 제안 샘플", 0.58, 0.78, 7.8, 1.65, 34, BLACK, True)
    add_text(slide, "295개 JSON을 API 없이 정제하고, 브랜드별 그룹과 추천 영업상품을 자동 산출한 예시입니다.", 0.62, 2.72, 6.7, 0.7, 16, MUTED)
    rect(slide, 8.1, 0.72, 4.55, 5.85, LIGHT)
    add_text(slide, "분석 흐름", 8.45, 1.05, 3.6, 0.42, 18, BLACK, True)
    steps = ["노이즈 제거", "리뷰/FAQ 문장 분리", "감성·기능 라벨링", "그룹/추천상품 매칭", "PPT·이미지 제안서 연결"]
    for i, s in enumerate(steps, start=1):
        y = 1.75 + (i - 1) * 0.78
        rect(slide, 8.45, y, 0.36, 0.36, ORANGE if i == 5 else BLACK)
        add_text(slide, f"{i}. {s}", 8.95, y - 0.03, 3.2, 0.36, 14, BLACK, True if i == 5 else False)
    add_footer(slide, 1)

    # 2. Snapshot
    slide = prs.slides.add_slide(blank)
    add_title(slide, "분석은 충분한 브랜드와 약한 브랜드를 먼저 갈라낸다", "근거등급은 제안 자동화에서 발송 우선순위를 정하는 기준이 됩니다.")
    metric(slide, "분석 브랜드", str(summary["brand_count"]), 0.58, 1.95, 2.25, BLACK)
    metric(slide, "A등급", str(summary["evidence_grade_counts"].get("A", 0)), 3.05, 1.95, 2.25, TEAL)
    metric(slide, "검토필요", str(summary["caution_counts"].get("상품군 자동 판정 불명확", 0)), 5.52, 1.95, 2.25, ORANGE)
    metric(slide, "대표문장 부족", str(summary["caution_counts"].get("대표 문장 부족", 0)), 7.99, 1.95, 2.25, BLUE)
    add_text(slide, "대표 그룹 분포", 0.58, 3.45, 3.2, 0.35, 18, BLACK, True)
    groups = list(summary["group_counts"].items())[:8]
    max_group = max(v for _, v in groups)
    for idx, (name, value) in enumerate(groups):
        bar(slide, name, value, max_group, 0.58, 4.02 + idx * 0.32, 6.4, TEAL if idx == 0 else BLACK)
    add_text(slide, "의미", 8.22, 3.45, 2.2, 0.35, 18, BLACK, True)
    bullet_list(slide, [
        "A/B등급은 바로 제안서 초안 생성 대상입니다.",
        "불명확 브랜드는 URL과 대표상품 확인 후 재분류합니다.",
        "그룹별 추천상품을 고정하면 대량 제안 메일로 확장하기 쉽습니다.",
    ], 8.22, 3.98, 4.25, 1.7, 14)
    add_footer(slide, 2)

    # 3. Recommended products
    slide = prs.slides.add_slide(blank)
    add_title(slide, "추천 영업상품은 상세페이지와 검색광고에 가장 많이 몰린다", "리뷰·FAQ에서 나온 구매장벽을 바로 제안상품으로 연결했습니다.")
    recs = list(summary["recommended_sales_product_counts"].items())
    max_rec = max(v for _, v in recs)
    for idx, (name, value) in enumerate(recs):
        y = 2.05 + idx * 0.58
        add_text(slide, name, 0.75, y - 0.07, 3.55, 0.35, 13, BLACK, True if idx < 2 else False)
        rect(slide, 4.55, y, 6.35, 0.22, RGBColor(232, 232, 232))
        rect(slide, 4.55, y, 6.35 * value / max_rec, 0.22, ORANGE if idx == 0 else BLACK)
        add_text(slide, str(value), 11.1, y - 0.08, 0.7, 0.32, 13, MUTED, align=PP_ALIGN.RIGHT)
    rect(slide, 0.74, 5.98, 11.85, 0.72, RGBColor(250, 242, 238))
    add_text(slide, "영업 메시지", 1.0, 6.13, 1.45, 0.28, 12, ORANGE, True)
    add_text(slide, "“리뷰와 FAQ에서 이미 드러난 구매장벽을 상세페이지·검색광고·숏폼 소재로 바꾸겠습니다.”", 2.42, 6.08, 9.6, 0.38, 16, BLACK, True)
    add_footer(slide, 3)

    # 4. Segment playbook
    slide = prs.slides.add_slide(blank)
    add_title(slide, "그룹별 제안은 상품이 아니라 문제 해결 단위로 묶는다", "같은 화장품이어도 고객이 말하는 불안과 기대가 다르면 제안상품도 달라집니다.")
    headers = ["그룹", "리뷰에서 잡는 신호", "추천 제안"]
    xs = [0.7, 3.5, 7.25]
    ws = [2.5, 3.25, 4.7]
    for x, w, h in zip(xs, ws, headers):
        rect(slide, x, 1.95, w, 0.45, BLACK)
        add_text(slide, h, x + 0.12, 2.04, w - 0.24, 0.22, 11, WHITE, True)
    playbook = [
        ("보습/수분", "건조, 당김, 촉촉, 흡수", "상세페이지 사용감 근거 강화"),
        ("장벽/진정", "시카, 판테놀, 자극 없음", "성분·효능 표현 리스크 검수"),
        ("트러블/피지", "여드름, 붉어짐, 모공", "검색광고 키워드와 전후 후기 소재"),
        ("선케어", "백탁, 눈시림, 발림성", "숏폼 사용감 테스트 콘텐츠"),
        ("가격/혜택형", "쿠폰, 무료배송, 재구매", "CRM·혜택 배너·리마케팅 소재"),
    ]
    for i, row in enumerate(playbook):
        y = 2.58 + i * 0.67
        for x, w, text in zip(xs, ws, row):
            rect(slide, x, y, w, 0.5, RGBColor(247, 247, 247) if i % 2 else WHITE, LINE)
            add_text(slide, text, x + 0.12, y + 0.1, w - 0.24, 0.24, 10, BLACK, True if x == xs[0] else False)
    add_footer(slide, 4)

    # 5. Brand sample
    slide = prs.slides.add_slide(blank)
    brand = sample["브랜드명"]
    add_title(slide, f"{brand} 예시: 리뷰 강점은 소재화하고, FAQ는 상세페이지로 흡수한다", "한 브랜드를 선택하면 자동으로 근거와 추천상품 브리프를 만들 수 있습니다.")
    metric(slide, "대표그룹", sample["대표그룹"], 0.7, 1.88, 3.0, TEAL)
    metric(slide, "세부유형", sample["세부유형"], 3.95, 1.88, 3.0, ORANGE)
    metric(slide, "근거등급", sample["근거등급"], 7.2, 1.88, 1.45, BLACK)
    metric(slide, "정제문장", sample["정제문장수"], 8.9, 1.88, 1.45, BLUE)
    add_text(slide, "대표 근거", 0.74, 3.35, 1.7, 0.32, 18, BLACK, True)
    bullet_list(slide, [clean_sentence(s, 76) for s in sample["대표문장"].split(" / ")[:2]], 0.78, 3.87, 5.5, 1.28, 12)
    add_text(slide, "추천 영업상품", 6.75, 3.35, 2.2, 0.32, 18, BLACK, True)
    bullet_list(slide, sample["추천영업상품"].split(" / ")[:4], 6.78, 3.87, 5.25, 1.25, 13)
    rect(slide, 6.78, 5.55, 5.25, 0.72, RGBColor(237, 247, 245))
    add_text(slide, "제안 문장", 7.02, 5.71, 1.25, 0.25, 11, TEAL, True)
    add_text(slide, clean_sentence(sample["제안브리프"], 72), 8.18, 5.66, 3.55, 0.34, 12, BLACK, True)
    add_footer(slide, 5)

    # 6. Close
    slide = prs.slides.add_slide(blank)
    add_text(slide, "다음 단계는 자동 제안서 생성으로 연결하는 것입니다", 0.58, 0.72, 10.8, 0.82, 30, BLACK, True)
    add_text(slide, "CSV의 그룹·근거·추천상품을 템플릿에 꽂으면 브랜드별 이미지와 PPT를 대량 생성할 수 있습니다.", 0.62, 1.68, 9.6, 0.45, 16, MUTED)
    actions = [
        ("1", "검토필요 74개 정리", "도메인과 대표상품을 확인해 오분류를 줄입니다."),
        ("2", "PPT 템플릿 고정", "영업상품별 3~5장 구조를 표준화합니다."),
        ("3", "이미지 카드 자동화", "대표 키워드와 추천상품을 썸네일 문구로 변환합니다."),
        ("4", "발송 우선순위", "A/B등급과 추천상품 수를 기준으로 메일 대상을 나눕니다."),
    ]
    for i, (num, head, body) in enumerate(actions):
        x = 0.72 + (i % 2) * 6.1
        y = 2.72 + (i // 2) * 1.45
        rect(slide, x, y, 5.55, 1.03, LIGHT)
        add_text(slide, num, x + 0.18, y + 0.18, 0.4, 0.38, 18, ORANGE, True)
        add_text(slide, head, x + 0.72, y + 0.16, 4.3, 0.3, 15, BLACK, True)
        add_text(slide, body, x + 0.72, y + 0.54, 4.55, 0.28, 10.5, MUTED)
    add_footer(slide, 6)

    prs.save(OUT_PPTX)
    print(OUT_PPTX)


if __name__ == "__main__":
    main()
