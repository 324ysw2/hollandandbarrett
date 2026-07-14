from __future__ import annotations

import argparse
import csv
import json
import re
from collections import Counter
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DATA_DIR = ROOT / "data" / "brand_json_summaries"
LOCAL_ANALYSIS_CSV = DATA_DIR / "ttamtti2_local_review_analysis.csv"
SEMANTIC_RECS_CSV = DATA_DIR / "ttamtti2_semantic_brand_recommendations.csv"
OUT_PREFIX = "ttamtti2_sales_automation_briefs"
ASSET_DIR = ROOT / "outputs" / "brand_sales_automation"

BLACK = RGBColor(15, 15, 15)
MUTED = RGBColor(86, 86, 86)
LIGHT = RGBColor(242, 242, 242)
LINE = RGBColor(190, 194, 200)
ORANGE = RGBColor(255, 107, 53)
TEAL = RGBColor(20, 126, 116)
BLUE = RGBColor(42, 91, 215)
PALE_TEAL = RGBColor(235, 247, 245)
PALE_ORANGE = RGBColor(255, 244, 239)
WHITE = RGBColor(255, 255, 255)


NOISE_WORDS = {
    "SHIPPING",
    "WORLD",
    "GUADELOUPE",
    "FRENCH",
    "REP",
    "USD",
    "KRW",
    "JPY",
    "Login",
    "Cart",
    "ABOUT",
    "Story",
}

NON_COSMETIC_TERMS = {
    "호텔",
    "맛집",
    "커피",
    "주식",
    "잠수함",
    "뉴스",
    "중공업",
    "게임",
    "메이플",
    "로아",
    "자동차",
    "시공",
    "랩핑",
}


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def by_brand(rows: list[dict[str, str]]) -> dict[str, dict[str, str]]:
    return {row["브랜드명"]: row for row in rows}


def split_slash(value: str, limit: int | None = None) -> list[str]:
    items = [" ".join(v.strip().split()) for v in (value or "").split(" / ") if v.strip()]
    return items[:limit] if limit else items


def split_comma(value: str, limit: int | None = None) -> list[str]:
    items = [" ".join(v.strip().split()) for v in (value or "").split(",") if v.strip()]
    items = [item for item in items if item and item not in NOISE_WORDS and item.upper() not in NOISE_WORDS]
    return items[:limit] if limit else items


def parse_aspects(value: str) -> dict[str, float]:
    out: dict[str, float] = {}
    for part in split_slash(value):
        if ":" not in part:
            continue
        key, raw = part.rsplit(":", 1)
        try:
            out[key.strip()] = float(raw.strip())
        except ValueError:
            pass
    return out


def to_int(value: str) -> int:
    try:
        return int(float(value or 0))
    except ValueError:
        return 0


def shorten(text: str, limit: int) -> str:
    text = re.sub(r"\([^)]*등록된 네이버 페이 구매평\)", " ", text or "")
    text = re.sub(r"\d{4}[-.]\d{1,2}[-.]\d{1,2}(?:\s+\d{1,2}:\d{2}:\d{2})?", " ", text)
    text = re.sub(r"[A-Za-z가-힣]\*{2,}", " ", text)
    text = " ".join(text.replace("더보기", "").split())
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "..."


def derive_barriers(local: dict[str, str], semantic: dict[str, str]) -> list[str]:
    aspects = parse_aspects(local.get("기능라벨분포", ""))
    barriers: list[tuple[str, int]] = []
    clean_count = to_int(local.get("정제문장수", "0"))
    pos_count = to_int(local.get("긍정문장수", "0"))
    neg_count = to_int(local.get("부정문장수", "0"))
    faq_count = to_int(local.get("FAQ문장수", "0"))
    grade = local.get("근거등급", "D")
    group = local.get("대표그룹", "")
    subtype = local.get("세부유형", "")
    recs = semantic.get("의미기반추천영업상품", "")
    caution = " ".join([local.get("주의플래그", ""), semantic.get("주의플래그", "")])

    if grade in {"C", "D"} or clean_count < 10 or "대표 문장 부족" in caution:
        barriers.append(("리뷰 근거 부족", 9))
    if faq_count >= max(2, pos_count + neg_count) or "FAQ/CS" in subtype:
        barriers.append(("FAQ/CS 정보 정리 필요", 8))
    if aspects.get("가격/혜택", 0) >= 2 or "가격/혜택" in subtype or local.get("가격신호"):
        barriers.append(("가격·혜택 설득 필요", 7))
    if aspects.get("배송/CS", 0) >= 1 or "배송" in local.get("부정주의키워드", ""):
        barriers.append(("배송·문의 불안 제거 필요", 7))
    if neg_count >= 2 or local.get("부정주의키워드"):
        barriers.append(("단점/주의 문구 보완 필요", 6))
    if "트러블" in group or "장벽" in group or aspects.get("트러블/자극", 0) >= 1:
        barriers.append(("자극·효능 신뢰 근거 필요", 6))
    if "화장품 표현 리스크 검수" in recs:
        barriers.append(("화장품 표현 리스크 확인 필요", 5))
    if aspects.get("사용감", 0) >= 1 or "선케어" in group or "클렌징" in group:
        barriers.append(("사용감 시각 증거 필요", 5))

    if not barriers:
        barriers.append(("강점 리뷰를 구매 설득 문구로 전환", 4))
    barriers.sort(key=lambda item: item[1], reverse=True)
    return [name for name, _ in barriers[:4]]


def derive_sentiment_summary(local: dict[str, str]) -> str:
    aspects = parse_aspects(local.get("기능라벨분포", ""))
    top = sorted(aspects.items(), key=lambda item: item[1], reverse=True)[:3]
    pos = split_comma(local.get("긍정키워드", ""), 4)
    neg = split_comma(local.get("부정주의키워드", ""), 3)
    top_text = ", ".join(name for name, _ in top) or "리뷰 항목 부족"
    pos_text = ", ".join(pos) if pos else "긍정 키워드 부족"
    neg_text = ", ".join(neg) if neg else "뚜렷한 부정 키워드 적음"
    return f"주요 항목은 {top_text}이며, 긍정 신호는 {pos_text}, 주의 신호는 {neg_text}입니다."


def priority_score(local: dict[str, str], semantic: dict[str, str], barriers: list[str]) -> int:
    score = 0
    score += {"A": 35, "B": 26, "C": 16, "D": 6}.get(local.get("근거등급", "D"), 10)
    score += min(to_int(local.get("정제문장수", "0")), 80) // 4
    if "리뷰/FAQ 기반 상세페이지 개선" in semantic.get("의미기반추천영업상품", ""):
        score += 12
    if "네이버 검색광고/쇼핑 키워드 패키지" in semantic.get("의미기반추천영업상품", ""):
        score += 8
    score += max(0, 16 - len(barriers) * 2)
    if "상품군 자동 판정 불명확" in local.get("주의플래그", ""):
        score -= 25
    joined = " ".join(
        [
            local.get("상위키워드", ""),
            local.get("대표문장", ""),
            semantic.get("공유키워드", ""),
            semantic.get("유사브랜드", ""),
        ]
    )
    if any(term in joined for term in NON_COSMETIC_TERMS):
        score -= 28
    return max(0, min(100, score))


def build_email(brand: str, semantic: dict[str, str], barriers: list[str], sentiment: str) -> tuple[str, str]:
    products = split_slash(semantic.get("의미기반추천영업상품", ""), 3)
    keywords = split_comma(semantic.get("공유키워드", ""), 5)
    similar = split_slash(semantic.get("유사브랜드", ""), 3)
    subject = f"{brand} 리뷰/FAQ 기준 제안 포인트 정리드립니다"
    body = "\n".join(
        [
            f"안녕하세요, {brand} 담당자님.",
            "",
            f"{brand}의 공개 리뷰와 FAQ 신호를 정리해보니 {', '.join(keywords[:3]) or '핵심 리뷰'} 키워드가 반복되고 있습니다.",
            f"유사한 반응을 보이는 브랜드로는 {', '.join(similar) or '동일 상품군 브랜드'}가 확인됩니다.",
            "",
            f"현재 가장 큰 구매장벽은 {', '.join(barriers[:2])}로 보이며, {sentiment}",
            "",
            "그래서 아래 제안을 우선 추천드립니다.",
            *[f"- {product}" for product in products],
            "",
            "리뷰 근거를 상세페이지와 광고 소재로 바꾸면 고객이 이미 말하고 있는 강점을 바로 구매 설득 문구로 전환할 수 있습니다.",
            "",
            "검토 가능하시면 브랜드별 예시 이미지와 1장 제안 PPT까지 함께 전달드리겠습니다.",
        ]
    )
    return subject, body


def sanitize_filename(name: str) -> str:
    cleaned = re.sub(r"[\\/:*?\"<>|]", "_", name)
    return cleaned.strip() or "brand"


def add_text(slide, text, x, y, w, h, size=16, color=BLACK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def rect(slide, x, y, w, h, fill=LIGHT, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def bullet_list(slide, items, x, y, w, h, size=11.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(7)
    return box


def make_ppt(brief: dict[str, Any], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    brand = brief["brand_name"]
    add_text(slide, f"{brand} 리뷰 기반 영업 제안", 0.62, 0.46, 8.8, 0.58, 30, BLACK, True)
    add_text(slide, "항목별 감성/구매장벽 분석과 유사 브랜드 추천을 연결한 자동 제안서 샘플입니다.", 0.66, 1.18, 10.2, 0.28, 12, MUTED)
    rect(slide, 0.66, 1.55, 11.95, 0.015, LINE)

    rect(slide, 0.72, 1.92, 2.75, 0.78, PALE_TEAL)
    add_text(slide, "그룹", 0.92, 2.05, 0.7, 0.18, 9, TEAL, True)
    add_text(slide, brief["main_group"], 0.92, 2.32, 2.25, 0.2, 12, BLACK, True)
    rect(slide, 3.72, 1.92, 2.45, 0.78, PALE_ORANGE)
    add_text(slide, "우선순위", 3.92, 2.05, 0.9, 0.18, 9, ORANGE, True)
    add_text(slide, f"{brief['priority_score']} / 100", 3.92, 2.32, 1.8, 0.2, 16, BLACK, True)

    add_text(slide, "구매장벽", 0.72, 3.12, 1.4, 0.25, 17, BLACK, True)
    bullet_list(slide, brief["purchase_barriers"], 0.78, 3.52, 5.25, 1.2, 12)

    add_text(slide, "추천 영업상품", 6.65, 3.12, 2.0, 0.25, 17, BLACK, True)
    bullet_list(slide, brief["recommended_sales_products"], 6.72, 3.52, 5.0, 1.2, 12)

    add_text(slide, "유사 브랜드", 0.72, 5.18, 1.4, 0.25, 17, BLACK, True)
    bullet_list(slide, brief["similar_brands"], 0.78, 5.56, 4.7, 0.8, 12)

    rect(slide, 6.65, 5.1, 5.32, 1.0, LIGHT)
    add_text(slide, "제안 메일 한 줄", 6.88, 5.28, 1.55, 0.2, 10.5, TEAL, True)
    add_text(slide, brief["email_one_liner"], 6.88, 5.6, 4.72, 0.28, 11.5, BLACK, True)

    add_text(slide, "Research-backed sales automation | review-based recommender + ABSA + proposal automation", 0.66, 7.08, 6.3, 0.18, 8, MUTED)
    prs.save(out_path)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        Path("C:/Windows/Fonts/malgunbd.ttf" if bold else "C:/Windows/Fonts/malgun.ttf"),
        Path("C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf"),
    ]
    for path in candidates:
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def wrap_text(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.ImageFont, max_width: int) -> list[str]:
    words = list(text) if re.search(r"[가-힣]", text) else text.split()
    lines: list[str] = []
    current = ""
    sep = "" if re.search(r"[가-힣]", text) else " "
    for word in words:
        candidate = current + (sep if current and sep else "") + word
        if draw.textbbox((0, 0), candidate, font=fnt)[2] <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def make_image(brief: dict[str, Any], out_path: Path) -> None:
    img = Image.new("RGB", (1200, 1500), "white")
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 1200, 1500), fill=(255, 255, 255))
    draw.rectangle((70, 80, 1130, 220), fill=(235, 247, 245))
    draw.text((100, 110), f"{brief['brand_name']} 리뷰 기반 제안", font=font(46, True), fill=(15, 15, 15))
    draw.text((100, 178), f"{brief['main_group']} | 우선순위 {brief['priority_score']}/100", font=font(24), fill=(20, 126, 116))

    y = 290
    sections = [
        ("구매장벽", brief["purchase_barriers"]),
        ("추천 영업상품", brief["recommended_sales_products"]),
        ("공유 키워드", brief["shared_keywords"][:7]),
        ("유사 브랜드", brief["similar_brands"][:4]),
    ]
    for title, items in sections:
        draw.text((90, y), title, font=font(30, True), fill=(255, 107, 53))
        y += 48
        for item in items:
            for line in wrap_text(draw, f"- {item}", font(25), 980)[:2]:
                draw.text((115, y), line, font=font(25), fill=(25, 25, 25))
                y += 34
            y += 8
        y += 34

    draw.rectangle((70, 1280, 1130, 1405), fill=(255, 244, 239))
    draw.text((100, 1312), "제안메일 한 줄", font=font(24, True), fill=(255, 107, 53))
    line_y = 1350
    for line in wrap_text(draw, brief["email_one_liner"], font(24, True), 930)[:2]:
        draw.text((100, line_y), line, font=font(24, True), fill=(15, 15, 15))
        line_y += 32
    img.save(out_path)


def build_brief(local: dict[str, str], semantic: dict[str, str]) -> dict[str, Any]:
    brand = local["브랜드명"]
    barriers = derive_barriers(local, semantic)
    sentiment = derive_sentiment_summary(local)
    subject, email_body = build_email(brand, semantic, barriers, sentiment)
    products = split_slash(semantic.get("의미기반추천영업상품", ""), 4)
    keywords = split_comma(semantic.get("공유키워드", ""), 10)
    similar = split_slash(semantic.get("유사브랜드", ""), 5)
    email_one_liner = (
        f"{brand}는 {', '.join(keywords[:3]) or '리뷰'} 신호가 있어 "
        f"{products[0] if products else '리뷰/FAQ 기반 상세페이지 개선'} 제안이 적합합니다."
    )
    return {
        "brand_name": brand,
        "main_group": local.get("대표그룹", ""),
        "subtype": local.get("세부유형", ""),
        "evidence_grade": local.get("근거등급", ""),
        "clean_sentence_count": to_int(local.get("정제문장수", "0")),
        "aspect_sentiment_summary": sentiment,
        "purchase_barriers": barriers,
        "shared_keywords": keywords,
        "similar_brands": similar,
        "recommended_sales_products": products,
        "recommendation_reason": semantic.get("추천근거", ""),
        "priority_score": priority_score(local, semantic, barriers),
        "email_subject": subject,
        "email_body": email_body,
        "email_one_liner": email_one_liner,
        "caution_flags": " / ".join(x for x in [local.get("주의플래그", ""), semantic.get("주의플래그", "")] if x),
    }


def write_outputs(briefs: list[dict[str, Any]], out_prefix: str) -> tuple[Path, Path, Path]:
    csv_path = DATA_DIR / f"{out_prefix}.csv"
    json_path = DATA_DIR / f"{out_prefix}.json"
    md_path = DATA_DIR / f"{out_prefix}.md"
    fields = [
        "브랜드명",
        "대표그룹",
        "세부유형",
        "근거등급",
        "우선순위점수",
        "항목별감성요약",
        "구매장벽",
        "공유키워드",
        "유사브랜드",
        "추천영업상품",
        "추천근거",
        "제안메일제목",
        "제안메일본문",
        "주의플래그",
    ]
    with csv_path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fields)
        writer.writeheader()
        for b in briefs:
            writer.writerow(
                {
                    "브랜드명": b["brand_name"],
                    "대표그룹": b["main_group"],
                    "세부유형": b["subtype"],
                    "근거등급": b["evidence_grade"],
                    "우선순위점수": b["priority_score"],
                    "항목별감성요약": b["aspect_sentiment_summary"],
                    "구매장벽": " / ".join(b["purchase_barriers"]),
                    "공유키워드": ", ".join(b["shared_keywords"]),
                    "유사브랜드": " / ".join(b["similar_brands"]),
                    "추천영업상품": " / ".join(b["recommended_sales_products"]),
                    "추천근거": b["recommendation_reason"],
                    "제안메일제목": b["email_subject"],
                    "제안메일본문": b["email_body"],
                    "주의플래그": b["caution_flags"],
                }
            )
    json_path.write_text(json.dumps(briefs, ensure_ascii=False, indent=2), encoding="utf-8")
    lines = [
        "# 브랜드 영업 자동화 브리프",
        "",
        "- 기반: review-based recommender systems, ABSA, B2B proposal automation",
        "- 흐름: 항목별 감성/구매장벽 분석 -> 유사 브랜드 기반 영업상품 추천 -> PPT/이미지/제안메일 자동 생성",
        "",
    ]
    for b in sorted(briefs, key=lambda x: x["priority_score"], reverse=True)[:30]:
        lines.extend(
            [
                f"## {b['brand_name']} ({b['priority_score']}/100)",
                f"- 감성/항목: {b['aspect_sentiment_summary']}",
                f"- 구매장벽: {' / '.join(b['purchase_barriers'])}",
                f"- 유사브랜드: {' / '.join(b['similar_brands'][:3])}",
                f"- 추천상품: {' / '.join(b['recommended_sales_products'])}",
                f"- 메일 한 줄: {b['email_one_liner']}",
                "",
            ]
        )
    md_path.write_text("\n".join(lines), encoding="utf-8")
    return csv_path, json_path, md_path


def select_asset_briefs(briefs: list[dict[str, Any]], brand: str | None, limit: int) -> list[dict[str, Any]]:
    if brand:
        for brief in briefs:
            if brief["brand_name"] == brand:
                return [brief]
        raise SystemExit(f"브랜드를 찾을 수 없습니다: {brand}")
    filtered = [
        b
        for b in briefs
        if not any(term in " ".join([b["brand_name"], ", ".join(b["shared_keywords"]), " / ".join(b["similar_brands"])]) for term in NON_COSMETIC_TERMS)
        and b["main_group"] != "검토필요/상품군 불명확"
    ]
    return sorted(filtered, key=lambda b: b["priority_score"], reverse=True)[:limit]


def write_assets(briefs: list[dict[str, Any]]) -> list[dict[str, str]]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    outputs: list[dict[str, str]] = []
    for brief in briefs:
        stem = sanitize_filename(brief["brand_name"])
        brand_dir = ASSET_DIR / stem
        brand_dir.mkdir(parents=True, exist_ok=True)
        ppt_path = brand_dir / f"{stem}_sales_brief.pptx"
        img_path = brand_dir / f"{stem}_sales_card.png"
        email_path = brand_dir / f"{stem}_proposal_email.txt"
        make_ppt(brief, ppt_path)
        make_image(brief, img_path)
        email_path.write_text(f"제목: {brief['email_subject']}\n\n{brief['email_body']}", encoding="utf-8")
        outputs.append({"brand": brief["brand_name"], "ppt": str(ppt_path), "image": str(img_path), "email": str(email_path)})
    return outputs


def main() -> int:
    parser = argparse.ArgumentParser(description="Build review-driven sales automation briefs and assets.")
    parser.add_argument("--brand", help="특정 브랜드만 PPT/이미지/메일 생성")
    parser.add_argument("--limit", type=int, default=3, help="브랜드 미지정 시 산출물 생성 개수")
    parser.add_argument("--no-assets", action="store_true", help="CSV/JSON/MD만 생성")
    args = parser.parse_args()

    local_rows = read_csv(LOCAL_ANALYSIS_CSV)
    semantic_rows = read_csv(SEMANTIC_RECS_CSV)
    semantic_by_brand = by_brand(semantic_rows)

    briefs = []
    for local in local_rows:
        semantic = semantic_by_brand.get(local["브랜드명"])
        if not semantic:
            continue
        briefs.append(build_brief(local, semantic))

    csv_path, json_path, md_path = write_outputs(briefs, OUT_PREFIX)
    asset_outputs: list[dict[str, str]] = []
    if not args.no_assets:
        asset_outputs = write_assets(select_asset_briefs(briefs, args.brand, args.limit))

    summary = {
        "brief_count": len(briefs),
        "csv": str(csv_path),
        "json": str(json_path),
        "markdown": str(md_path),
        "assets": asset_outputs,
    }
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
